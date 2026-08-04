"""
IA Processor para EducaApp
===========================
Módulo actualizado con PyMuPDF y DocumentProcessor para procesamiento avanzado.
Mantiene retrocompatibilidad con funciones existentes.
"""

# from transformers import pipeline
import fitz  # PyMuPDF (reemplaza a PyPDF2)
from docx import Document
from pptx import Presentation
import os
import re
import sys

# Importar el DocumentProcessor del módulo nuevo
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from document_processor import DocumentProcessor


# Instancia global del procesador para reutilizar
_processor = DocumentProcessor()


def extract_text_from_file(file_path):
    """
    Extrae texto de archivos PDF, DOCX, PPTX, TXT.
    ACTUALIZADO: Usa PyMuPDF en vez de PyPDF2 para PDFs.
    
    Args:
        file_path: Ruta al archivo
    
    Returns:
        Texto extraído del documento
    """
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()
    text = ""

    if file_extension == '.pdf':
        # Usar PyMuPDF (fitz) en vez de PyPDF2
        doc = fitz.open(file_path)
        for i, page in enumerate(doc):
            page_text = page.get_text()
            if page_text:
                text += f"\n[Página {i+1}]\n" + page_text
        doc.close()
        
    elif file_extension == '.docx':
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
            
    elif file_extension == '.pptx':
        ppt = Presentation(file_path)
        for i, slide in enumerate(ppt.slides, 1):
            text += f"\n[Slide {i}]\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                    
    elif file_extension == '.txt':
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
            
    else:
        raise ValueError(f"Formato de archivo no soportado: {file_extension}")
    
    return text.strip()


def extract_page_images(file_path, pages=None, max_images=4, min_bytes=4000):
    """
    Extrae imágenes embebidas de un PDF (diagramas, fotos, gráficos) usando
    PyMuPDF, para mandarlas a un modelo de IA con visión junto con el texto.
    Por ahora solo soporta PDF — DOCX/PPTX devuelven lista vacía (no es la
    fuente más común para este caso de uso; se puede sumar después si hace
    falta).

    Args:
        file_path: ruta al archivo.
        pages: lista opcional de números de página (1-based) a considerar;
            None = todo el documento.
        max_images: tope de imágenes a devolver (controla costo/latencia de
            la llamada a la IA — cada imagen se manda entera en el prompt).
        min_bytes: descarta imágenes muy chicas (íconos, viñetas, líneas
            decorativas del layout) que no aportan contenido real.

    Returns:
        Lista de dicts {page, mime, data_uri}, como mucho `max_images`,
        ordenada por página. Nunca levanta excepción — ante cualquier error
        de lectura devuelve lista vacía (la generación de preguntas sigue
        funcionando solo con texto).
    """
    import base64

    if os.path.splitext(file_path)[1].lower() != '.pdf':
        return []

    results = []
    try:
        doc = fitz.open(file_path)
        try:
            page_range = range(len(doc)) if not pages else [p - 1 for p in pages if 0 < p <= len(doc)]
            for page_idx in page_range:
                if len(results) >= max_images:
                    break
                page = doc[page_idx]
                for img in page.get_images(full=True):
                    if len(results) >= max_images:
                        break
                    xref = img[0]
                    try:
                        extracted = doc.extract_image(xref)
                    except Exception:
                        continue
                    image_bytes = extracted.get('image')
                    ext = (extracted.get('ext') or 'png').lower()
                    if not image_bytes or len(image_bytes) < min_bytes:
                        continue
                    if ext not in ('png', 'jpg', 'jpeg', 'webp'):
                        continue
                    mime = 'image/jpeg' if ext in ('jpg', 'jpeg') else f'image/{ext}'
                    b64 = base64.b64encode(image_bytes).decode('utf-8')
                    results.append({
                        'page': page_idx + 1,
                        'mime': mime,
                        'data_uri': f'data:{mime};base64,{b64}',
                    })
        finally:
            doc.close()
    except Exception:
        return []

    return results


def extract_text_advanced(file_path, remove_headers=True, remove_footers=True):
    """
    NUEVA FUNCIÓN: Extracción avanzada con limpieza de headers/footers.
    Usa el DocumentProcessor completo.
    
    Args:
        file_path: Ruta al archivo
        remove_headers: Eliminar headers repetitivos (PDFs)
        remove_footers: Eliminar footers repetitivos (PDFs)
    
    Returns:
        Dict con estructura completa del documento
    """
    from django.conf import settings

    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    if file_extension == '.pdf':
        return _processor.process_pdf(
            file_path,
            remove_headers=remove_headers,
            remove_footers=remove_footers,
            extract_toc=True,
            max_pages=getattr(settings, 'CONTENIDO_MAX_PAGES', None),
            pages_per_block=getattr(settings, 'CONTENIDO_PAGES_PER_CHAPTER_BLOCK', 40),
        )
    elif file_extension == '.docx':
        return _processor.process_docx(file_path)
    elif file_extension == '.pptx':
        return _processor.process_pptx(file_path)
    else:
        # Fallback para otros formatos
        text = extract_text_from_file(file_path)
        return {
            'metadata': {'title': os.path.basename(file_path)},
            'chapters': [{
                'title': 'Documento completo',
                'content': text,
                'tokens': _processor.count_tokens(text)
            }],
            'stats': {'total_tokens': _processor.count_tokens(text)}
        }


def count_tokens(text):
    """
    NUEVA FUNCIÓN: Cuenta tokens exactos usando tiktoken (GPT-4).
    
    Args:
        text: Texto a contar
    
    Returns:
        Número de tokens
    """
    return _processor.count_tokens(text)


def count_tokens_file(file_path):
    """
    NUEVA FUNCIÓN: Cuenta tokens de un archivo completo.
    
    Args:
        file_path: Ruta al archivo
    
    Returns:
        Número total de tokens
    """
    result = extract_text_advanced(file_path)
    return result['stats'].get('total_tokens', 0)


def split_text_by_tokens(text, max_tokens=4000):
    """
    NUEVA FUNCIÓN: Divide texto en chunks que no excedan límite de tokens.
    
    Args:
        text: Texto a dividir
        max_tokens: Límite de tokens por chunk
    
    Returns:
        Lista de chunks de texto
    """
    return _processor.split_by_token_limit(text, max_tokens=max_tokens)


def optimize_text_for_ai(text, remove_extra_whitespace=True):
    """
    NUEVA FUNCIÓN: Optimiza texto para reducir tokens sin perder información.
    
    Args:
        text: Texto a optimizar
        remove_extra_whitespace: Elimina espacios/saltos extras
    
    Returns:
        Texto optimizado
    """
    return _processor.optimize_for_ai(
        text,
        remove_extra_whitespace=remove_extra_whitespace,
        remove_urls=False,  # Por defecto no quitar URLs en documentos académicos
        remove_emails=False
    )


def split_text_into_chapters(text):
    """
    Divide el texto en capítulos basados en títulos detectados.
    NOTA: Para mejor detección, usar extract_text_advanced() que usa TOC.
    """
    chapters = re.split(r'(?i)(cap[ií]tulo \d+|secci[oó]n \d+)', text)
    return [c.strip() for c in chapters if c.strip()]


def extract_book_metadata(file_path):
    """
    Extrae metadata de un libro PDF (ISBN, título, autor, edición, editorial, año, páginas).
    
    Args:
        file_path: Ruta al archivo PDF
    
    Returns:
        dict con metadata extraída: {
            'title': str,
            'author': str,
            'isbn': str,
            'edition': str,
            'publisher': str,
            'year': int,
            'pages': int
        }
    """
    metadata = {
        'title': '',
        'author': '',
        'isbn': '',
        'edition': '',
        'publisher': '',
        'year': None,
        'pages': None
    }
    
    _, file_extension = os.path.splitext(file_path)
    if file_extension.lower() != '.pdf':
        return metadata
    
    try:
        doc = fitz.open(file_path)
        
        # 1. Metadata del PDF (título, autor, fecha)
        pdf_metadata = doc.metadata
        if pdf_metadata:
            metadata['title'] = pdf_metadata.get('title', '').strip()
            metadata['author'] = pdf_metadata.get('author', '').strip()
            
            # Intentar extraer año de creationDate o modDate
            creation_date = pdf_metadata.get('creationDate', '')
            if creation_date and len(creation_date) >= 4:
                try:
                    # Formato: D:YYYYMMDD...
                    year_match = re.search(r'(\d{4})', creation_date)
                    if year_match:
                        metadata['year'] = int(year_match.group(1))
                except:
                    pass
        
        # 2. Número de páginas
        metadata['pages'] = len(doc)
        
        # 3. Extraer texto de las primeras 10 páginas para buscar ISBN, edición, editorial
        # (muchos libros tienen la info de copyright en páginas 3-10)
        first_pages_text = ""
        for page_num in range(min(10, len(doc))):
            first_pages_text += f"\n--- Página {page_num+1} ---\n"
            first_pages_text += doc[page_num].get_text()
        
        # 4. Buscar ISBN (varios formatos, incluyendo con guiones)
        isbn_patterns = [
            r'ISBN[:\s-]*(\d{3}[-\s]?\d{1,5}[-\s]?\d{1,7}[-\s]?\d{1,7}[-\s]?\d)',  # ISBN-13 con guiones
            r'ISBN[:\s-]*(\d{13})',  # ISBN-13 sin guiones
            r'ISBN[:\s-]*(\d{1,5}[-\s]?\d{1,7}[-\s]?\d{1,7}[-\s]?[\dX])',  # ISBN-10 con guiones
            r'ISBN[:\s-]*(\d{10})',  # ISBN-10 sin guiones
            r'(?:ISBN|isbn)[-:]?\s*(\d{10,13})',  # ISBN simple
        ]
        
        for pattern in isbn_patterns:
            isbn_match = re.search(pattern, first_pages_text, re.IGNORECASE)
            if isbn_match:
                isbn = isbn_match.group(1)
                # Normalizar: mantener solo dígitos (eliminar guiones y espacios)
                isbn_clean = isbn.replace('-', '').replace(' ', '')
                # Validar longitud (10 o 13 dígitos)
                if len(isbn_clean) in [10, 13]:
                    metadata['isbn'] = isbn  # Mantener formato original con guiones
                    break
        
        # 5. Buscar edición (patrones mejorados para español)
        edition_patterns = [
            r'Primera\s+edici[oó]n[:\s]+(\w+\s+de\s+)?(\d{4})',  # "Primera edición: noviembre de 2024"
            r'(\d+)[ºª]?\s+edici[oó]n',  # "3ª edición" o "3 edición"
            r'(\d+)(?:st|nd|rd|th)\s+edition',  # "3rd edition"
            r'edici[oó]n[:\s]+(\d+)',  # "Edición: 3"
            r'edition[:\s]+(\d+)',  # "Edition: 3"
        ]
        
        for pattern in edition_patterns:
            edition_match = re.search(pattern, first_pages_text, re.IGNORECASE)
            if edition_match:
                # Para "Primera edición", extraer 1
                if 'primera' in pattern.lower() or 'first' in pattern.lower():
                    metadata['edition'] = '1'
                else:
                    # Extraer el número de edición
                    groups = edition_match.groups()
                    for group in groups:
                        if group and group.isdigit():
                            metadata['edition'] = group
                            break
                if metadata['edition']:
                    break
        
        # 6. Buscar editorial (patrones mejorados, ordenados por prioridad)
        publisher_patterns = [
            # Primero: patrones más específicos
            r'Ediciones\s+([A-Z][A-ZÁ-ÿ]+)',  # "Ediciones OCTAEDRO"
            r'Editorial\s+([A-Z][A-Za-zÁ-ÿ]+(?:\s+[A-Z][A-Za-zÁ-ÿ]+)?)',  # "Editorial Pearson Education"
            # Segundo: con dos puntos
            r'Editorial[:\s]+([A-Z][A-Za-zÁ-ÿ\s&,.-]+?)(?:\n|,)',  # "Editorial: Nombre"
            r'Publisher[:\s]+([A-Z][A-Za-z\s&,.-]+?)(?:\n|,)',  # "Publisher: Name"
            r'Publicado\s+por[:\s]+([A-Z][A-Za-zÁ-ÿ\s&,.-]+?)(?:\n|,)',  # "Publicado por: Editorial"
            r'Published\s+by[:\s]+([A-Z][A-Za-z\s&,.-]+?)(?:\n|,)',  # "Published by: Publisher"
        ]
        
        for pattern in publisher_patterns:
            publisher_match = re.search(pattern, first_pages_text)
            if publisher_match:
                publisher = publisher_match.group(1).strip()
                # Limpiar publisher (quitar espacios extras, comas finales, etc.)
                publisher = re.sub(r'\s+', ' ', publisher)
                publisher = publisher.strip(',.-')
                # Validar longitud razonable
                if 2 < len(publisher) < 80:
                    metadata['publisher'] = publisher
                    break
        
        # 7. Buscar año de publicación en el texto (patrones mejorados)
        if not metadata['year']:
            year_patterns = [
                r'Primera\s+edici[oó]n[:\s]+\w+\s+de\s+(\d{4})',  # "Primera edición: noviembre de 2024"
                r'©\s*(\d{4})',  # "© 2024"
                r'Copyright\s*©?\s*(\d{4})',  # "Copyright © 2024"
                r'Publicado\s+en\s+(\d{4})',  # "Publicado en 2024"
                r'Published\s+in\s+(\d{4})',  # "Published in 2024"
                r'(\d{4})\s*©',  # "2024 ©"
            ]
            
            for pattern in year_patterns:
                year_match = re.search(pattern, first_pages_text, re.IGNORECASE)
                if year_match:
                    year = int(year_match.group(1))
                    if 1900 <= year <= 2030:  # Validar rango razonable
                        metadata['year'] = year
                        break
        
        # 8. Si no hay título en metadata, intentar extraer de la primera página
        if not metadata['title']:
            lines = first_pages_text.split('\n')
            # Buscar la primera línea con longitud significativa (probablemente el título)
            for line in lines[:10]:
                line = line.strip()
                if 10 < len(line) < 200 and not line.startswith('ISBN'):
                    metadata['title'] = line
                    break
        
        doc.close()
        
    except Exception as e:
        print(f"Error extrayendo metadata: {e}")
    
    return metadata
