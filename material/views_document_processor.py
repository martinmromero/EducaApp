"""
Vistas para Document Processing en EducaApp
============================================
Endpoints para procesar documentos, contar tokens y preparar contenido para IA.
"""

from django.shortcuts import render, get_object_or_404
from django.http import JsonResponse, StreamingHttpResponse, FileResponse, Http404
from django.views.decorators.http import require_http_methods
from django.contrib.auth.decorators import login_required
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
from django.conf import settings
import os
import re
import uuid
import shutil
import tempfile
import threading
import time
import logging
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# In-memory job store for SSE streaming jobs (job_id → params)
_jobs = {}
_jobs_lock = threading.Lock()

from material.ia_processor import (
    extract_text_advanced,
    extract_page_images,
    count_tokens,
    split_text_by_tokens,
    optimize_text_for_ai
)
from material.local_ai_client import local_ai

logger = logging.getLogger(__name__)

# document_processor.py incrusta marcadores invisibles de página
# (\x00P<pagina física>\x00, ver _join_pages_with_markers) en el 'content'
# de cada capítulo, para que _split_into_chunks pueda citar de qué página
# salió cada fragmento realmente mandado a la IA (ver Question.source_page /
# source_chapters más abajo). Nunca deben llegar a algo que se le muestra al
# usuario (content_preview) — quedarían como "P23" pegado en medio del texto.
_PAGE_MARKER_RE = re.compile(r'\x00P(\d+)\x00')


def _strip_page_markers(text):
    return _PAGE_MARKER_RE.sub('', text or '')


@login_required
@require_http_methods(["POST"])
def upload_and_process_document(request):
    """
    Vista para subir y procesar un documento (PDF, DOCX, PPTX).
    Retorna estructura completa con capítulos, tokens, metadata.
    También guarda el archivo como Contenido en la base de datos.
    
    POST params:
        - documento: archivo subido
        - contenido_title: título para Mis Contenidos (opcional)
        - remove_headers: bool (opcional, default True)
        - remove_footers: bool (opcional, default True)
    
    Returns:
        JSON con estructura procesada del documento
    """
    from material.models import Contenido, Subject

    if 'documento' not in request.FILES:
        return JsonResponse({
            'success': False,
            'error': 'No se envió ningún archivo'
        }, status=400)
    
    archivo = request.FILES['documento']
    remove_headers = request.POST.get('remove_headers', 'true').lower() == 'true'
    remove_footers = request.POST.get('remove_footers', 'true').lower() == 'true'
    contenido_title = request.POST.get('contenido_title', '').strip()
    subject_id = request.POST.get('subject_id', '').strip()

    # Validar tamaño (plan gratuito de Render: memoria y tiempo de request limitados)
    max_upload_mb = settings.CONTENIDO_MAX_UPLOAD_MB
    if archivo.size > max_upload_mb * 1024 * 1024:
        return JsonResponse({
            'success': False,
            'error': (
                f'El archivo pesa {archivo.size / (1024 * 1024):.1f}MB y supera el máximo permitido '
                f'de {max_upload_mb}MB. Si es un PDF escaneado, probar comprimirlo o subir solo las páginas necesarias.'
            )
        }, status=400)

    # Validar extensión
    nombre = archivo.name
    ext = os.path.splitext(nombre)[1].lower()
    if ext not in ['.pdf', '.docx', '.pptx', '.txt']:
        return JsonResponse({
            'success': False,
            'error': f'Formato no soportado: {ext}'
        }, status=400)

    if not contenido_title:
        contenido_title = os.path.splitext(nombre)[0].replace('_', ' ').replace('-', ' ')

    try:
        import hashlib
        import tempfile
        from .cleanup import compute_file_hash

        # Leer los bytes del archivo UNA sola vez (el stream no es re-readable)
        file_bytes = archivo.read()
        file_hash = compute_file_hash(file_bytes)

        # ---- Deduplicación ----
        existing = Contenido.objects.filter(
            file_hash=file_hash, uploaded_by=request.user
        ).first()

        duplicate_message = None

        if existing and existing.file_available and existing.file_actually_exists():
            # Archivo idéntico ya existe y sigue vigente — no guardamos un nuevo archivo
            contenido_id = existing.id
            duplicate_message = (
                f'Este documento ya existe como "{existing.title}". '
                f'Se usará el archivo guardado; no se creó un duplicado.'
            )
            # Procesar con archivo temporal (compatible con local y cloud)
            ext = os.path.splitext(nombre)[1].lower()
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                result = extract_text_advanced(
                    tmp_path,
                    remove_headers=remove_headers,
                    remove_footers=remove_footers
                )
            finally:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass
            file_path = tmp_path  # Solo para la sesión; apunta al archivo original en el caso disponible
            # Para la sesión usamos el path del archivo existente si es local
            try:
                session_file_path = existing.file.path
            except (ValueError, NotImplementedError, AttributeError):
                session_file_path = tmp_path  # fallback cloud

        elif existing:
            # Archivo había expirado, o el registro decía "vigente" pero el archivo
            # físico ya no está (borrado por reinicio/limpieza sin actualizar
            # file_deleted_at) — en ambos casos, restaurarlo con el nuevo upload.
            saved_relative = default_storage.save(f'contenidos/{nombre}', ContentFile(file_bytes))
            file_path = os.path.join(settings.MEDIA_ROOT, saved_relative)
            existing.file = saved_relative
            existing.file_deleted_at = None
            existing.save(update_fields=['file', 'file_deleted_at'])
            contenido_id = existing.id
            duplicate_message = (
                f'El archivo de "{existing.title}" había expirado y fue restaurado correctamente.'
            )
            result = extract_text_advanced(
                file_path,
                remove_headers=remove_headers,
                remove_footers=remove_footers
            )
            session_file_path = file_path

        else:
            # Documento nuevo
            saved_relative = default_storage.save(f'contenidos/{nombre}', ContentFile(file_bytes))
            file_path = os.path.join(settings.MEDIA_ROOT, saved_relative)

            result = extract_text_advanced(
                file_path,
                remove_headers=remove_headers,
                remove_footers=remove_footers
            )

            contenido = Contenido(
                title=contenido_title,
                uploaded_by=request.user,
                file_hash=file_hash,
            )
            contenido.file = saved_relative
            contenido.save()
            contenido_id = contenido.id
            session_file_path = file_path
        # Etiquetar la materia elegida (si vino informada) sobre el Contenido resultante,
        # sea nuevo, restaurado o reusado por deduplicación.
        if subject_id.isdigit():
            try:
                subj = Subject.objects.get(pk=int(subject_id))
                Contenido.objects.get(pk=contenido_id).subjects.add(subj)
            except Subject.DoesNotExist:
                pass

        # --- Actualizar sesión ---
        # Eliminar archivo previo de session temporal si era una sesión de doc_sessions
        prev_session = request.session.get('doc_processor', {})
        prev_path = prev_session.get('file_path', '')
        sessions_dir = os.path.join(settings.MEDIA_ROOT, 'doc_sessions')
        if prev_path and prev_path.startswith(str(sessions_dir)) and os.path.exists(prev_path):
            try:
                os.unlink(prev_path)
            except OSError:
                pass

        doc_id = str(uuid.uuid4())
        request.session['doc_processor'] = {
            'doc_id': doc_id,
            'file_path': session_file_path,
            'filename': nombre,
            'remove_headers': remove_headers,
            'remove_footers': remove_footers,
        }
        request.session.modified = True
        # ---------------------------------------------------

        total_tokens = result.get('stats', {}).get('total_tokens', 0)

        # Formatear respuesta (content_preview solo para mostrar en UI)
        response_data = {
            'success': True,
            'doc_id': doc_id,
            'filename': nombre,
            'contenido_id': contenido_id,
            'duplicate_message': duplicate_message,
            'metadata': result.get('metadata', {}),
            'stats': result.get('stats', {}),
            'scanned_pages': result.get('scanned_pages', []),
            'chapters': [
                {
                    'title': ch.get('title', ''),
                    'tokens': ch.get('tokens', 0),
                    'content_preview': _strip_page_markers(ch.get('content', ''))[:6000],
                    'pages': ch.get('pages', [])
                }
                for ch in result.get('chapters', [])
            ],
            'toc': result.get('toc', []),
            # Presupuestos de tokens, para que la UI oriente la selección de capítulos
            'token_budget': {
                'total_budget': settings.CONTENIDO_MAX_TOTAL_TOKENS,
                'run_budget': settings.CONTENIDO_MAX_RUN_TOKENS,
                'exceeds_total_budget': total_tokens > settings.CONTENIDO_MAX_TOTAL_TOKENS,
            },
        }

        return JsonResponse(response_data)
        
    except Exception as e:
        # Limpiar en caso de error
        if 'file_path' in locals() and os.path.exists(file_path):
            try:
                os.unlink(file_path)
            except OSError:
                pass
        
        return JsonResponse({
            'success': False,
            'error': str(e)
        }, status=500)


@login_required
@require_http_methods(["POST"])
def split_text_chunks(request):
    """
    Divide un texto en chunks según límite de tokens.
    
    POST params:
        - text: texto a dividir
        - max_tokens: límite por chunk (default: 4000)
    
    Returns:
        JSON con chunks
    """
    text = request.POST.get('text', '')
    max_tokens = int(request.POST.get('max_tokens', 4000))
    
    if not text:
        return JsonResponse({
            'success': False,
            'error': 'No se proporcionó texto'
        }, status=400)
    
    try:
        chunks = split_text_by_tokens(text, max_tokens=max_tokens)
        
        return JsonResponse({
            'success': True,
            'total_chunks': len(chunks),
            'max_tokens_per_chunk': max_tokens,
            'chunks': [
                {
                    'chunk_number': i + 1,
                    'tokens': count_tokens(chunk),
                    'preview': chunk[:100] + '...' if len(chunk) > 100 else chunk
                }
                for i, chunk in enumerate(chunks)
            ]
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e)
        }, status=500)


@login_required
def document_processor_dashboard(request):
    """
    Vista HTML para el dashboard de procesamiento de documentos.
    """
    from material.content_visibility import get_visible_subjects
    # Obtener backend configurado para este usuario (no la instancia global)
    from .ai_router import get_backend_for_user, get_global_demo_quota, ensure_fresh_demo_quota, SharedDemoBackend
    backend = get_backend_for_user(request.user)
    ai_status = backend.get_status()

    using_shared_fallback = isinstance(backend, SharedDemoBackend)
    if using_shared_fallback:
        ensure_fresh_demo_quota()
    demo_quota = get_global_demo_quota() if using_shared_fallback else None
    demo_quota_low_tokens = bool(
        demo_quota and demo_quota.get('limit_tokens')
        and demo_quota.get('remaining_tokens') is not None
        and (demo_quota['remaining_tokens'] / demo_quota['limit_tokens']) < 0.2
    )

    # ONBOARDING WIZARD V2: si venimos del asistente (?wizard=1), lo recordamos en
    # sesión para poder mostrar el banner de continuidad también en /crear-examen/
    # más adelante, sin tener que pasar el parámetro a mano por todos lados.
    # Si en cambio se entra por la navegación normal (sidebar "Procesador IA",
    # "Mis Contenidos"), sin el parámetro, limpiamos cualquier flag viejo para
    # que no quede un banner de "seguís en el asistente" pegado de una sesión
    # anterior que se abandonó a mitad de camino.
    if request.GET.get('wizard') == '1':
        request.session['onb2_wizard_active'] = True
    else:
        request.session.pop('onb2_wizard_active', None)
    wizard_active = request.session.get('onb2_wizard_active', False)

    preselected_subject_id = request.GET.get('subject_id', '')
    preselected_subject_name = ''
    if preselected_subject_id.isdigit():
        preselected_subject_name = (
            get_visible_subjects(request.user)
            .filter(id=preselected_subject_id)
            .values_list('name', flat=True)
            .first() or ''
        )

    context = {
        'page_title': 'Procesador de Documentos',
        'supported_formats': ['.pdf', '.docx', '.pptx', '.txt'],
        'max_file_size_mb': settings.CONTENIDO_MAX_UPLOAD_MB,
        'total_token_budget': settings.CONTENIDO_MAX_TOTAL_TOKENS,
        'run_token_budget': settings.CONTENIDO_MAX_RUN_TOKENS,
        'local_ai_connected': ai_status.get('connected', False),
        'local_ai_ready': ai_status.get('ready_for_generation', ai_status.get('connected', False)),
        'selected_model': ai_status.get('selected_model', ai_status.get('model', 'N/A')),
        'default_model': ai_status.get('default_model', ai_status.get('model', 'N/A')),
        'backend_type': ai_status.get('backend', 'ollama_local'),
        'preselected_contenido_id': request.GET.get('contenido_id', ''),
        'preselected_subject_id': preselected_subject_id,
        'preselected_subject_name': preselected_subject_name,
        'wizard_active': wizard_active,
        'using_shared_fallback': using_shared_fallback,
        'demo_quota': demo_quota,
        'demo_quota_low_tokens': demo_quota_low_tokens,
    }
    
    return render(request, 'material/document_processor_dashboard.html', context)


@login_required
def process_contenido_by_id(request, contenido_id):
    """
    Procesa un Contenido ya guardado en el servidor y devuelve el mismo
    formato JSON que upload_and_process_document, para preseleccionarlo
    en el dashboard del doc-processor.
    """
    from material.models import Contenido
    from django.http import JsonResponse

    contenido = get_object_or_404(Contenido, id=contenido_id, uploaded_by=request.user)

    if not contenido.file or not contenido.file_actually_exists():
        return JsonResponse({
            'success': False,
            'error': f'El archivo de "{contenido.title}" ya no está disponible en el servidor '
                     f'(se elimina automáticamente después de unos días). Subir el documento de nuevo.',
        }, status=404)

    file_path = contenido.file.path
    ext = os.path.splitext(file_path)[1].lower()

    if ext not in ['.pdf', '.docx', '.pptx', '.txt']:
        return JsonResponse({'success': False, 'error': f'Formato no soportado: {ext}'}, status=400)

    try:
        result = extract_text_advanced(file_path, remove_headers=True, remove_footers=True)

        # Limpiar sesion previa
        prev_session = request.session.get('doc_processor', {})
        prev_path = prev_session.get('file_path')
        if prev_path and os.path.exists(prev_path) and not prev_path.startswith(str(settings.MEDIA_ROOT).rstrip('/') + '/contenidos'):
            try:
                os.unlink(prev_path)
            except OSError:
                pass

        doc_id = str(uuid.uuid4())
        request.session['doc_processor'] = {
            'doc_id': doc_id,
            'file_path': file_path,
            'filename': contenido.file.name.split('/')[-1],
            'remove_headers': True,
            'remove_footers': True,
        }
        request.session.modified = True

        nombre = contenido.file.name.split('/')[-1]
        return JsonResponse({
            'success': True,
            'doc_id': doc_id,
            'filename': nombre,
            'contenido_id': contenido.id,
            'metadata': result.get('metadata', {}),
            'stats': result.get('stats', {}),
            'scanned_pages': result.get('scanned_pages', []),
            'chapters': [
                {
                    'title': ch.get('title', ''),
                    'tokens': ch.get('tokens', 0),
                    'content_preview': _strip_page_markers(ch.get('content', ''))[:6000],
                    'pages': ch.get('pages', [])
                }
                for ch in result.get('chapters', [])
            ],
            'toc': result.get('toc', []),
            'token_budget': {
                'total_budget': settings.CONTENIDO_MAX_TOTAL_TOKENS,
                'run_budget': settings.CONTENIDO_MAX_RUN_TOKENS,
                'exceeds_total_budget': result.get('stats', {}).get('total_tokens', 0) > settings.CONTENIDO_MAX_TOTAL_TOKENS,
            },
        })

    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


# ============================================
# Vistas para Servidor Local de IA
# ============================================

@login_required
@require_http_methods(["GET"])
def check_local_ai_status(request):
    """
    Verifica el estado de conexión al proveedor de IA configurado por el usuario.

    Returns:
        JSON con estado de conexión y modelo activo
    """
    try:
        from .ai_router import get_backend_for_user, get_global_demo_quota, ensure_fresh_demo_quota, SharedDemoBackend
        from .models import UserAIConfig
        config, _ = UserAIConfig.objects.get_or_create(user=request.user)
        backend = get_backend_for_user(request.user)
        status = backend.get_status()
        # OJO: NO pisar status['backend'] con config.source. config.source es
        # la preferencia guardada por el usuario ('ollama_local', 'byok', etc.),
        # pero get_backend_for_user() puede resolver a un backend real distinto
        # (ej. Ollama no disponible → cae al fallback compartido de Groq) sin
        # tocar esa preferencia. Pisar 'backend' acá rompía el chequeo de "Sin
        # límite por tanda" (piensa que sigue en Ollama local cuando en
        # realidad está usando el fallback compartido con cupo limitado).
        status['source'] = config.source
        if isinstance(backend, SharedDemoBackend):
            status['using_shared_fallback'] = True
            ensure_fresh_demo_quota()
            quota = get_global_demo_quota()
            if quota:
                status['demo_quota'] = {
                    'provider': quota['provider'],
                    'remaining_requests': quota['remaining_requests'],
                    'limit_requests': quota['limit_requests'],
                    'requests_reset_at': quota['requests_reset_at'].isoformat() if quota['requests_reset_at'] else None,
                    'remaining_tokens': quota['remaining_tokens'],
                    'limit_tokens': quota['limit_tokens'],
                }
        return JsonResponse({'success': True, **status})
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e),
            'connected': False
        }, status=500)


@login_required
@require_http_methods(["GET"])
def list_local_ai_models(request):
    """
    Lista todos los modelos disponibles en el servidor local.
    
    Returns:
        JSON con lista de modelos
    """
    try:
        models = local_ai.get_models()
        current_model = local_ai.get_current_model()
        
        return JsonResponse({
            'success': True,
            'connected': local_ai.is_available(),
            'models': models,
            'count': len(models),
            'selected_model': current_model,
            'default_model': local_ai.default_model
        })
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e),
            'models': []
        }, status=500)


@login_required
@require_http_methods(["POST"])
def set_local_ai_model(request):
    """
    Cambia el modelo activo del servidor local.
    
    POST params:
        - model: nombre del modelo a activar
    
    Returns:
        JSON con resultado de la operación
    """
    try:
        model_name = request.POST.get('model')
        
        if not model_name:
            return JsonResponse({
                'success': False,
                'error': 'Nombre de modelo no proporcionado'
            }, status=400)
        
        # Intentar cambiar el modelo
        success = local_ai.set_model(model_name)
        
        if success:
            return JsonResponse({
                'success': True,
                'message': f'Modelo cambiado a {model_name}',
                'selected_model': model_name
            })
        else:
            return JsonResponse({
                'success': False,
                'error': 'Modelo no disponible o servidor no conectado'
            }, status=400)
            
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e)
        }, status=500)


@login_required
@require_http_methods(["POST"])
def generate_questions_from_chapters(request):
    """
    Genera preguntas automáticamente usando IA desde capítulos seleccionados.

    Lee el documento completo persistido en sesión y aplica chunking para
    procesar capítulos de cualquier longitud sin truncar el contenido.

    POST JSON:
        - chapter_indices: lista de índices de capítulos (preferido)
        - chapters: lista de capítulos con título (fallback si no hay sesión)
        - filename: nombre del archivo fuente
        - doc_id: ID de documento (opcional, para validación)

    Returns:
        JSON con preguntas generadas
    """
    import json as json_module

    try:
        data = json_module.loads(request.body)
        chapter_indices = data.get('chapter_indices', [])
        chapters_from_request = data.get('chapters', [])
        filename = data.get('filename', 'Documento')
        stream_mode = data.get('stream_mode', False)

        if not chapter_indices and not chapters_from_request:
            return JsonResponse({
                'success': False,
                'error': 'No se proporcionaron capítulos'
            }, status=400)

        # Verificar servidor IA
        from .ai_router import get_backend_for_user
        _ai_backend = get_backend_for_user(request.user)
        _status = _ai_backend.get_status()
        if not _status.get('connected'):
            backend_type = _status.get('backend', 'ollama_local')
            if backend_type == 'ollama_local':
                error_msg = (
                    'Servidor Ollama no disponible. '
                    'Para usar el generador de IA en producción, configurar un proveedor '
                    'en "Proveedor de IA" (BYOK con OpenAI, Anthropic, etc.).'
                )
            else:
                error_msg = f'Proveedor de IA ({backend_type}) no disponible. Verificar la configuración y la API key.'
            return JsonResponse({
                'success': False,
                'error': error_msg
            }, status=503)

        _content_chunk_tokens, _output_tokens_ceiling = _chunking_budget(_ai_backend)

        # question_types enviados por el cliente (lista de strings)
        question_types = data.get('question_types') or []

        # Cantidad total de preguntas deseadas (usuario elige, default 20, máximo 200)
        total_questions = max(1, min(200, int(data.get('total_questions', 20) or 20)))

        # --------------------------------------------------------
        # Preguntas ya existentes para este documento (anti-repetición)
        # --------------------------------------------------------
        contenido_id = data.get('contenido_id')
        existing_texts_set = set()
        existing_questions_list = []
        if contenido_id:
            try:
                existing_texts_set, existing_questions_list = _get_existing_questions_for_contenido(
                    contenido_id, request.user
                )
                if existing_questions_list:
                    logger.info(f"Contenido {contenido_id}: {len(existing_questions_list)} preguntas previas en BD → incluidas en prompt anti-repetición")
            except Exception as _e:
                logger.warning(f"No se pudo obtener preguntas existentes: {_e}")

        # Modo streaming: guardar job y retornar job_id al cliente
        if stream_mode:
            # 0 = "calcular automáticamente" (ver stream_questions): max(1, ...)
            # rompía ese sentinel forzando siempre 1 pregunta por fragmento, sin
            # importar total_questions/total_chunks. max(0, ...) preserva el 0.
            questions_per_block = max(0, min(12, int(data.get('questions_per_block', 0) or 0)))
            job_id = _store_streaming_job(
                request, chapter_indices, chapters_from_request, filename,
                question_types, total_questions, questions_per_block,
                existing_questions_list=existing_questions_list,
                existing_texts_set=existing_texts_set,
                include_images=bool(data.get('include_images')),
            )
            return JsonResponse({
                'success': True,
                'job_id': job_id,
                'stream': True,
                'existing_count': len(existing_questions_list),
            })

        # --------------------------------------------------------
        # Obtener contenido completo desde el archivo en sesión
        # --------------------------------------------------------
        doc_session = request.session.get('doc_processor', {})
        session_file = doc_session.get('file_path')
        chapters_to_process = []

        if session_file and os.path.exists(session_file):
            # Re-procesar el archivo original para obtener contenido completo
            try:
                full_result = extract_text_advanced(
                    session_file,
                    remove_headers=doc_session.get('remove_headers', True),
                    remove_footers=doc_session.get('remove_footers', True)
                )
                all_session_chapters = full_result.get('chapters', [])

                if chapter_indices:
                    chapters_to_process = [
                        all_session_chapters[i]
                        for i in chapter_indices
                        if i < len(all_session_chapters)
                    ]
                else:
                    # Hacer match por título con los capítulos del request
                    request_titles = {ch.get('title', '') for ch in chapters_from_request}
                    chapters_to_process = [
                        ch for ch in all_session_chapters
                        if ch.get('title', '') in request_titles
                    ]
            except Exception as e:
                logger.warning(f"No se pudo re-procesar el archivo de sesión: {e}")
                # Caer en modo fallback

        # Fallback: usar el contenido preview que vino en el request
        if not chapters_to_process:
            logger.info("Usando contenido del request como fallback (sin sesión de archivo)")
            chapters_to_process = chapters_from_request

        if not chapters_to_process:
            return JsonResponse({
                'success': False,
                'error': 'No se pudo obtener el contenido de los capítulos'
            }, status=400)

        selected_tokens = _chapters_total_tokens(chapters_to_process)
        run_budget = settings.CONTENIDO_MAX_RUN_TOKENS
        if selected_tokens > run_budget:
            return JsonResponse({
                'success': False,
                'error': (
                    f'Se seleccionaron {_fmt_es(selected_tokens)} tokens de contenido, y el máximo por '
                    f'tanda es {_fmt_es(run_budget)}. Elegir menos capítulos y generar el resto en otra tanda.'
                ),
            }, status=400)

        # --------------------------------------------------------
        # Generar preguntas por capítulo usando chunking
        # --------------------------------------------------------
        all_questions = []
        failed_chunks = []

        # Total de chunks de TODOS los capítulos seleccionados, calculado antes de
        # generar nada — total_questions se distribuye contra este total global,
        # no contra los chunks de cada capítulo por separado (eso hacía que pedir
        # 20 preguntas con 3 capítulos terminara generando ~20 por capítulo, unas
        # 60 en total).
        chapter_chunks = [
            (chapter, _split_into_chunks(
                chapter.get('content', chapter.get('content_preview', '')), max_tokens=_content_chunk_tokens
            ))
            for chapter in chapters_to_process
        ]
        total_chunks_all = sum(len(chunks) for _, chunks in chapter_chunks)

        for chapter, chunks in chapter_chunks:
            title = chapter.get('title', 'Capítulo')
            chapter_pages = chapter.get('pages', [])
            # Física → impresa (ver document_processor._detect_printed_page_number):
            # si el libro tiene su propia numeración (ej. un capítulo aislado
            # que en el PDF empieza en la página física 2 pero dice "322"
            # impreso), se cita esa — si no se detectó, se cae a la física.
            printed_map = dict(zip(chapter_pages, chapter.get('printed_pages', [])))
            logger.info(f"Procesando capítulo '{title}' ({len(chunks)} chunk(s))")

            # Piso de 1 (no 2): con muchos chunks chicos, un piso más alto
            # infla el total muy por encima de lo pedido. Techo de 12: pedir más
            # por fragmento excede lo que el max_tokens de salida (4096) puede
            # sostener, y Groq puede directamente rechazar la request (413).
            questions_per_chunk = max(1, min(12, total_questions // max(total_chunks_all, 1)))

            for chunk_idx, chunk in enumerate(chunks):
                if not _chunk_has_content(chunk['text']):
                    logger.warning(f"Chunk {chunk_idx + 1}/{len(chunks)} de '{title}' sin texto suficiente, se omite.")
                    failed_chunks.append({
                        'chapter': title, 'chunk': chunk_idx + 1, 'total_chunks': len(chunks),
                        'error': 'Fragmento sin texto extraíble (posible página escaneada o solo con imágenes) — no se generaron preguntas.',
                    })
                    continue
                # Pequeño respiro entre requests: pedir varios fragmentos
                # seguidos sin pausa puede superar el límite por minuto (TPM/RPM)
                # de Groq incluso con cada request individual dentro de lo
                # permitido — el backend ya reintenta ante un 429, pero es mejor
                # no provocarlo de entrada.
                if chunk_idx > 0:
                    time.sleep(2)
                try:
                    chunk_questions = _generate_questions_for_chunk(
                        chunk['text'], title, questions_per_chunk, chunk_idx, len(chunks),
                        question_types=question_types, backend=_ai_backend,
                        existing_questions=existing_questions_list,
                        output_tokens_ceiling=_output_tokens_ceiling,
                    )
                except Exception as exc:
                    logger.warning(f"Chunk {chunk_idx + 1}/{len(chunks)} de '{title}' falló: {exc}")
                    failed_chunks.append({'chapter': title, 'chunk': chunk_idx + 1, 'total_chunks': len(chunks), 'error': str(exc)})
                    continue
                # Tope duro: el modelo no siempre respeta "generá exactamente N
                # preguntas" al pie de la letra, puede devolver bastantes más.
                remaining = max(0, total_questions - len(all_questions))
                chunk_questions = chunk_questions[:remaining]
                # Cada pregunta se etiqueta con las páginas del FRAGMENTO puntual
                # del que realmente salió (no con todo el capítulo/tanda).
                chunk_pages = chunk['pages'] or chapter_pages
                display_pages = sorted({printed_map.get(p) or p for p in chunk_pages}) if chunk_pages else []
                for q in chunk_questions:
                    q['source_chapters'] = [{'title': title, 'pages': display_pages}]
                    q['source_file'] = filename
                all_questions.extend(chunk_questions)
                logger.info(f"  chunk {chunk_idx + 1}/{len(chunks)}: {len(chunk_questions)} preguntas")

                if len(all_questions) >= total_questions:
                    break

            if len(all_questions) >= total_questions:
                break

        # Deduplicar contra preguntas ya en BD y entre sí
        all_questions = _deduplicate_questions(all_questions, extra_seen=existing_texts_set)

        if not all_questions:
            backend_status = _ai_backend.get_status() if _ai_backend else {}
            error_msg = (
                f'La IA no pudo generar preguntas válidas en {len(failed_chunks)} de los fragmentos procesados.'
                if failed_chunks else
                'La IA respondió, pero no generó preguntas válidas. '
                'Revisar el proveedor y el modelo configurado, especialmente en Gemini.'
            )
            return JsonResponse({
                'success': False,
                'error': error_msg,
                'failed_chunks': failed_chunks,
                'backend': backend_status.get('backend', 'unknown'),
                'provider': backend_status.get('provider', ''),
                'model': backend_status.get('model', ''),
            }, status=422)

        logger.info(f"Total preguntas generadas (dedup): {len(all_questions)}")

        return JsonResponse({
            'success': True,
            'questions': all_questions,
            'count': len(all_questions),
            'existing_count': len(existing_questions_list),
            'failed_chunks': failed_chunks,
        })

    except Exception as e:
        logger.exception("Error en generate_questions_from_chapters")
        return JsonResponse({
            'success': False,
            'error': str(e)
        }, status=500)


def _store_streaming_job(request, chapter_indices, chapters_from_request, filename,
                         question_types=None, total_questions=20, questions_per_block=0,
                         existing_questions_list=None, existing_texts_set=None,
                         include_images=False):
    """Guarda los parámetros del job en memoria y retorna el job_id."""
    job_id = str(uuid.uuid4())
    with _jobs_lock:
        # Limpiar jobs viejos (> 10 min) para no acumular basura
        now = time.time()
        stale = [jid for jid, j in _jobs.items() if now - j.get('created_at', 0) > 600]
        for jid in stale:
            _jobs.pop(jid, None)

        _jobs[job_id] = {
            'chapter_indices': chapter_indices,
            'chapters_from_request': chapters_from_request,
            'filename': filename,
            'question_types': question_types or [],
            'total_questions': total_questions,
            'questions_per_block': questions_per_block,  # 0 = calcular automáticamente
            'doc_session': dict(request.session.get('doc_processor', {})),
            'user_id': request.user.id,
            'created_at': now,
            'existing_questions_list': existing_questions_list or [],
            'existing_texts_set': existing_texts_set or set(),
            'include_images': include_images,
        }
    return job_id


# ============================================
# Helpers para generación por chunks
# ============================================

def _chapters_total_tokens(chapters):
    """Suma los tokens (ya calculados en la extracción) de una lista de capítulos."""
    return sum(ch.get('tokens', 0) for ch in chapters)


def _fmt_es(n):
    """Formatea un entero con punto como separador de miles (convención local)."""
    return f'{n:,}'.replace(',', '.')


# Por debajo de esto, un chunk no tiene texto real para generar preguntas
# (páginas escaneadas sin OCR, slides solo con imágenes, secciones vacías
# entre encabezados) — sin este piso, el prompt le llegaba a la IA con
# "TEXTO:" seguido de nada, y el modelo generaba preguntas sobre el propio
# prompt (ej. "¿qué tipo de pregunta se indica como 'opcion_multiple'?")
# en vez de sobre el documento. Ver [[project_fotosintesis_prompt_leak]].
MIN_CHUNK_CHARS = 40


def _chunk_has_content(chunk):
    return bool(chunk and len(chunk.strip()) >= MIN_CHUNK_CHARS)


# Defaults históricos: 3000 tokens de contenido por fragmento, hasta 4096 de
# salida. Le quedan cortos a algunos modelos gratuitos del fallback
# compartido (ej. llama-3.1-8b-instant en Groq, con un límite de apenas 6000
# tokens por minuto/TPM): contenido + resto del prompt + salida reservada
# supera ese TPM y Groq rechaza la request entera con "Request too large"
# — no es un tema de reintentos, la request en sí ya viene sobredimensionada.
_DEFAULT_CONTENT_CHUNK_TOKENS = 3000
_DEFAULT_OUTPUT_TOKENS_CEILING = 4096


def _chunking_budget(backend):
    """Cuánto contenido mandar por fragmento y cuántos tokens de salida pedir,
    respetando el TPM real del proveedor cuando se conoce (fallback
    compartido de demo, vía el último cupo capturado de sus headers reales —
    ver GlobalFallbackBackend._save_quota_snapshot). Con proveedores propios
    (BYOK) no se conoce ese límite de antemano, así que se usan los defaults
    de siempre.

    Devuelve (content_max_tokens, output_max_tokens_ceiling).
    """
    from .ai_router import SharedDemoBackend, get_global_demo_quota

    if not isinstance(backend, SharedDemoBackend):
        return _DEFAULT_CONTENT_CHUNK_TOKENS, _DEFAULT_OUTPUT_TOKENS_CEILING

    quota = get_global_demo_quota()
    tpm = quota.get('limit_tokens') if quota else None
    if not tpm:
        return _DEFAULT_CONTENT_CHUNK_TOKENS, _DEFAULT_OUTPUT_TOKENS_CEILING

    # Reparto conservador del TPM real: ~35% contenido, ~50% salida, ~15% de
    # margen para el resto del prompt (instrucciones, preguntas ya generadas
    # para no repetir, etc.) que no se cuenta acá.
    content_budget = max(600, min(_DEFAULT_CONTENT_CHUNK_TOKENS, int(tpm * 0.35)))
    output_budget = max(400, min(_DEFAULT_OUTPUT_TOKENS_CEILING, int(tpm * 0.5)))
    return content_budget, output_budget


def _split_into_chunks(content, max_tokens=3000):
    """Divide el contenido en fragmentos de ≤ max_tokens tokens.

    Devuelve una lista de dicts {'text': ..., 'pages': [...]} — 'text' es el
    texto del fragmento listo para mandar a la IA (sin marcadores), y
    'pages' son las páginas físicas del PDF de origen que aportaron texto a
    ESE fragmento puntual (no las de todo el capítulo), reconstruidas a
    partir de los marcadores invisibles que document_processor.py incrusta
    en 'content' (ver _join_pages_with_markers). Para contenido sin
    marcadores (DOCX/PPTX/TXT, que no tienen ese concepto de página física),
    'pages' queda vacía — quien llama debe caer de vuelta a las páginas del
    capítulo entero en ese caso.
    """
    raw_chunks = _split_into_chunks_raw(content, max_tokens=max_tokens)
    result = []
    for raw in raw_chunks:
        pages = sorted({int(p) for p in _PAGE_MARKER_RE.findall(raw)})
        result.append({'text': _PAGE_MARKER_RE.sub('', raw).strip(), 'pages': pages})
    return result


def _split_into_chunks_raw(content, max_tokens=3000):
    """Algoritmo de fragmentado en sí (por párrafos, y por líneas si un
    párrafo solo ya supera max_tokens) — ver _split_into_chunks (el que hay
    que usar desde afuera) para la envoltura consciente de páginas."""
    total_tokens = count_tokens(content)
    if total_tokens <= max_tokens:
        return [content]

    # Dividir por párrafos y agrupar hasta el límite
    paragraphs = content.split('\n\n')
    chunks = []
    current_parts = []
    current_tokens = 0

    for para in paragraphs:
        para_tokens = count_tokens(para)
        # Si un párrafo solo ya supera el límite, cortarlo por líneas
        if para_tokens > max_tokens:
            if current_parts:
                chunks.append('\n\n'.join(current_parts))
                current_parts = []
                current_tokens = 0
            lines = para.split('\n')
            for line in lines:
                line_tokens = count_tokens(line)
                if current_tokens + line_tokens > max_tokens and current_parts:
                    chunks.append('\n'.join(current_parts))
                    current_parts = [line]
                    current_tokens = line_tokens
                else:
                    current_parts.append(line)
                    current_tokens += line_tokens
        elif current_tokens + para_tokens > max_tokens and current_parts:
            chunks.append('\n\n'.join(current_parts))
            current_parts = [para]
            current_tokens = para_tokens
        else:
            current_parts.append(para)
            current_tokens += para_tokens

    if current_parts:
        chunks.append('\n\n'.join(current_parts))

    return chunks if chunks else [content]


def _build_generation_prompt(context):
    """
    Arma el prompt real a partir del template guardado en QuestionGenerationConfig
    (editable desde Administración → "Prompt de generación IA"). Si no hay fila
    todavía, o el template guardado tiene un placeholder inválido/typo que rompe
    el .format(), cae al default de fábrica — la generación de preguntas nunca
    debería fallar por un prompt mal editado.

    Devuelve (prompt_text, temperature).
    """
    from .models import QuestionGenerationConfig
    from .ai_prompts import DEFAULT_PROMPT_TEMPLATE, DEFAULT_TEMPERATURE

    cfg = QuestionGenerationConfig.objects.first()
    template = cfg.prompt_template if cfg and cfg.prompt_template else DEFAULT_PROMPT_TEMPLATE
    temperature = cfg.temperature if cfg else DEFAULT_TEMPERATURE

    try:
        return template.format(**context), temperature
    except (KeyError, ValueError, IndexError) as e:
        logger.error(
            f"Prompt guardado en QuestionGenerationConfig tiene un placeholder inválido "
            f"({e}) — usando el default de fábrica para esta generación."
        )
        return DEFAULT_PROMPT_TEMPLATE.format(**context), DEFAULT_TEMPERATURE


def _generate_questions_for_chunk(content, chapter_title, num_questions, chunk_idx, total_chunks, question_types=None, backend=None, existing_questions=None, images=None, output_tokens_ceiling=None, generate_kwargs=None):
    """Genera preguntas para un fragmento de capítulo usando la IA configurada.

    Args:
        content: Texto del fragmento a procesar.
        chapter_title: Título del capítulo.
        num_questions: Número total de preguntas a generar.
        chunk_idx: Índice del fragmento actual (0-based).
        total_chunks: Total de fragmentos del capítulo.
        question_types: Lista de tipos habilitados.
        existing_questions: lista de dicts {pregunta, respuesta, tipo} ya guardadas en BD.
        images: lista opcional de data-URIs (ver ia_processor.extract_page_images) para
            mandar junto con el texto a un modelo con visión. Si el modelo configurado
            no soporta imágenes, la llamada falla y este chunk se reporta como error
            (igual que cualquier otro fallo de chunk) — es responsabilidad de quien
            llama no pasar `images` salvo que el usuario lo haya pedido explícitamente.
        output_tokens_ceiling: tope de tokens de salida a pedirle al modelo (ver
            _chunking_budget) — None usa el default histórico (4096).
        generate_kwargs: dict opcional de kwargs extra para backend.generate() —
            ej. reasoning_effort para modelos de razonamiento (gpt-oss, qwen3.x
            en Groq). Los backends que no reconocen el kwarg lo ignoran (todos
            aceptan **kwargs), así que es seguro pasarlo sin importar el backend.
    """
    import json as json_module

    # Tipos disponibles y sus descripciones para el prompt
    ALL_TYPES = {
        'opcion_multiple': 'Opción múltiple (4 opciones A/B/C/D, una correcta)',
        'verdadero_falso': 'Verdadero/Falso (afirmación clara)',
        'completar_blank': 'Completar el espacio (usa [___] en la pregunta para el espacio en blanco)',
        'desarrollo': 'Desarrollo (pregunta abierta con respuesta de referencia para el docente)',
    }

    if not question_types:
        question_types = list(ALL_TYPES.keys())

    enabled_descriptions = '\n'.join(
        f'  - "{t}": {ALL_TYPES[t]}'
        for t in question_types
        if t in ALL_TYPES
    )

    context_note = f"(parte {chunk_idx + 1} de {total_chunks})" if total_chunks > 1 else ""

    # Niveles de Bloom para el prompt
    bloom_desc = (
        "bloom_nivel: nivel cognitivo de Bloom (1=Recordar, 2=Comprender, 3=Aplicar, "
        "4=Analizar, 5=Evaluar, 6=Crear)"
    )

    # Bloque de preguntas ya existentes para evitar repeticiones
    existing_block = ""
    if existing_questions:
        # Limitar a 40 para no inflar el prompt innecesariamente
        sample = existing_questions[:40]
        lines = [f'  {i+1}. [{q["tipo"]}] {q["pregunta"]}' for i, q in enumerate(sample)]
        existing_block = (
            f"\n\nPREGUNTAS YA GENERADAS PARA ESTE DOCUMENTO (NO REPETIR NI PARAFRASEAR):\n"
            + "\n".join(lines)
            + "\n\nGenerá preguntas completamente distintas a las anteriores, sobre aspectos o ángulos diferentes del texto.\n"
        )

    if images and not (content and content.strip()):
        # Página escaneada sin capa de texto: no hay "TEXTO:" real más abajo,
        # así que hay que decírselo explícito para que no invente contenido
        # ni pregunte sobre la ausencia de texto — ver [[project_fotosintesis_prompt_leak]].
        images_note = (
            "\nEste fragmento no tiene texto extraíble (es una página escaneada). "
            "Se incluyen las imágenes de esas páginas: generá las preguntas basándote "
            "pura y exclusivamente en lo que se ve en ellas.\n"
        )
    elif images:
        images_note = (
            "\nTambién se incluyen una o más imágenes de este fragmento del documento "
            "(diagramas, gráficos, fotos). Si aportan contenido educativo relevante, "
            "generá al menos una pregunta que haga referencia a lo que se ve en ellas.\n"
        )
    else:
        images_note = ""

    prompt_context = {
        'num_questions': num_questions,
        'chapter_title': chapter_title,
        'context_note': context_note,
        'images_note': images_note,
        'content': content,
        'enabled_descriptions': enabled_descriptions,
        'existing_block': existing_block,
        'bloom_desc': bloom_desc,
    }
    prompt, temperature = _build_generation_prompt(prompt_context)

    # ~300 tokens por pregunta (opcion_multiple/desarrollo con explicación suelen ser
    # las más largas) + margen fijo. Antes esto era un 4000 fijo sin importar cuántas
    # preguntas se pedían: con más de ~12-15 preguntas por chunk, la respuesta se
    # cortaba a mitad de camino y el parseo de JSON fallaba o rescataba solo 1-2.
    # Techo bajado de 8192 a 4096: pedirle 8192 a Groq (ej. documentos cortos con
    # pocos fragmentos, donde num_questions termina siendo alto por chunk) dio
    # "413 Payload Too Large" — el proveedor rechaza la request directamente en
    # vez de truncar. 4096 es más conservador; si se pide más de ~12 preguntas
    # en un mismo chunk, igual puede no alcanzar y quedar corto, pero no falla.
    #
    # Modelos de razonamiento (gpt-oss, qwen3.x en Groq) gastan tokens del
    # mismo presupuesto en pensar antes de escribir el JSON — la estimación
    # de arriba no deja margen para eso, así que se le suma un colchón fijo.
    # Default a RESERVARLO (no a 0 salvo que se pida explícito): este call
    # site casi nunca sabe de antemano qué modelo va a atender la llamada
    # (ai_router.GROQ_REASONING_MODEL_DEFAULTS ya completa reasoning_effort
    # solo al construir el payload, y DemoRoutingBackend recién decide
    # Groq/Gemini adentro de generate()) — antes, al no pasar generate_kwargs
    # acá, este colchón quedaba en 0 pese a que el modelo sí razonaba con el
    # esfuerzo default de Groq, dejando a veces la respuesta completamente
    # vacía (json_validate_failed con failed_generation="", ver GroqMonitorRun
    # 2026-08-27). Pedir de más no cuesta nada: se factura y corta por uso
    # real, no por el techo pedido — por eso solo se saca el colchón cuando
    # quien llama sabe con certeza que no hace falta (reasoning_effort="none").
    reasoning_effort = (generate_kwargs or {}).get('reasoning_effort')
    reasoning_budget = 0 if reasoning_effort == 'none' else 2000
    gen_max_tokens = min(
        output_tokens_ceiling or _DEFAULT_OUTPUT_TOKENS_CEILING,
        300 * max(num_questions, 1) + 500 + reasoning_budget,
    )

    # Ollama (local_ai/OllamaBackend) no tiene parámetro `images` — solo se lo
    # pasamos al backend externo, y solo cuando hay imágenes de verdad, para
    # no romper la firma de generate() de ningún backend que no lo espere.
    extra_kwargs = {'images': [img['data_uri'] for img in images]} if images else {}
    # Pedir JSON estructurado a nivel API (no solo por instrucción de prompt) en
    # los backends que lo soportan (OpenAICompatibleBackend — ver ai_router.py);
    # los demás lo ignoran vía **kwargs.
    extra_kwargs['json_mode'] = True
    if generate_kwargs:
        extra_kwargs.update(generate_kwargs)
    if backend is not None:
        result = backend.generate(prompt=prompt, temperature=temperature, max_tokens=gen_max_tokens, **extra_kwargs)
    else:
        result = local_ai.generate(prompt=prompt, temperature=temperature, max_tokens=gen_max_tokens)

    if not result['success']:
        error_msg = result.get('error', 'Error desconocido del proveedor de IA')
        logger.warning(f"IA falló para chunk {chunk_idx + 1}: {error_msg}")
        raise RuntimeError(error_msg)

    if result.get('truncated'):
        logger.warning(
            f"Respuesta de IA truncada por límite de tokens en chunk {chunk_idx + 1}/{total_chunks} "
            f"de '{chapter_title}' (max_tokens={gen_max_tokens}, {num_questions} preguntas pedidas). "
            "Es posible que se recuperen menos preguntas de las pedidas."
        )

    try:
        ai_response = result['text'].strip()
        # Eliminar bloques de código markdown si existen
        if ai_response.startswith('```'):
            lines = ai_response.split('\n')
            ai_response = '\n'.join(line for line in lines if not line.startswith('```'))
        # Intentar extraer el JSON si viene con texto alrededor
        start = ai_response.find('{')
        end = ai_response.rfind('}')
        if start != -1 and end != -1 and end > start:
            ai_response = ai_response[start:end + 1]
        questions_data = json_module.loads(ai_response)
        questions = questions_data.get('preguntas', [])
        # Filtrar solo los tipos habilitados (la IA puede equivocarse)
        questions = [q for q in questions if q.get('tipo', 'opcion_multiple') in question_types]
        return questions
    except Exception as e:
        logger.warning(f"No se pudo parsear JSON del chunk {chunk_idx + 1}: {e}")
        truncation_note = ' (truncado por límite de tokens)' if result.get('truncated') else ''
        raise RuntimeError(
            f'La IA respondió, pero el contenido no tenía el formato esperado (fragmento {chunk_idx + 1} de {total_chunks} de "{chapter_title}"){truncation_note}.'
        ) from e


def _first_source_page(source_chapters):
    """Extrae la primera página detectada en source_chapters, para poblar
    Question.source_page (mismo campo que usa la carga manual)."""
    for chapter in source_chapters or []:
        pages = chapter.get('pages') or []
        if pages:
            return pages[0]
    return None


def _deduplicate_questions(questions, extra_seen=None):
    """Elimina duplicados comparando los primeros 80 chars (case-insensitive).

    extra_seen: set adicional de claves (primeros 120 chars) ya vistas en BD.
    """
    seen = set(extra_seen) if extra_seen else set()
    unique = []
    for q in questions:
        key80  = q.get('pregunta', '').lower().strip()[:80]
        key120 = q.get('pregunta', '').lower().strip()[:120]
        if key80 and key80 not in seen and key120 not in seen:
            seen.add(key80)
            unique.append(q)
    return unique


def _get_existing_questions_for_contenido(contenido_id, user):
    """
    Devuelve (texts_set, summary_list) con todas las preguntas ya guardadas
    en la BD que apuntan a este Contenido (sin importar estado IA).

    texts_set   → set de str (primeros 120 chars en minúscula) para dedup rápido.
    summary_list → lista de dicts {pregunta, respuesta, tipo} para incluir en el prompt.
    """
    from material.models import Question
    qs = Question.objects.filter(
        contenido_id=contenido_id,
        user=user,
    ).values('question_text', 'answer_text', 'question_type')

    texts_set = set()
    summary_list = []
    for row in qs:
        txt = row['question_text'].strip()
        key = txt.lower()[:120]
        texts_set.add(key)
        summary_list.append({
            'pregunta': txt,
            'respuesta': (row['answer_text'] or '').strip(),
            'tipo': row['question_type'],
        })
    return texts_set, summary_list


@login_required
@require_http_methods(["GET"])
def stream_questions(request, job_id):
    """
    SSE endpoint: transmite preguntas a medida que se generan por chunks.
    El cliente abre un EventSource hacia esta URL tras recibir job_id del
    endpoint POST /generate-questions/ con stream_mode=true.
    """
    import json as json_module

    def event_stream(job, user_id, backend):
        # Verificar usuario antes de procesar
        if job.get('user_id') != user_id:
            yield f'data: {json_module.dumps({"type": "error", "message": "No autorizado"})}\n\n'
            return

        chapter_indices = job['chapter_indices']
        chapters_from_request = job['chapters_from_request']
        filename = job['filename']
        doc_session = job['doc_session']
        question_types = job.get('question_types') or []
        total_questions = max(1, int(job.get('total_questions', 20) or 20))
        questions_per_block_override = int(job.get('questions_per_block', 0) or 0)
        existing_questions_list = job.get('existing_questions_list') or []
        existing_texts_set = set(job.get('existing_texts_set') or [])
        include_images = bool(job.get('include_images'))
        content_chunk_tokens, output_tokens_ceiling = _chunking_budget(backend)

        # Obtener contenido completo (misma lógica que generate_questions_from_chapters)
        chapters_to_process = []
        session_file = doc_session.get('file_path')
        if session_file and os.path.exists(session_file):
            try:
                full_result = extract_text_advanced(
                    session_file,
                    remove_headers=doc_session.get('remove_headers', True),
                    remove_footers=doc_session.get('remove_footers', True)
                )
                all_session_chapters = full_result.get('chapters', [])
                if chapter_indices:
                    chapters_to_process = [
                        all_session_chapters[i]
                        for i in chapter_indices
                        if i < len(all_session_chapters)
                    ]
                else:
                    req_titles = {ch.get('title', '') for ch in chapters_from_request}
                    chapters_to_process = [
                        ch for ch in all_session_chapters
                        if ch.get('title', '') in req_titles
                    ]
            except Exception as exc:
                logger.warning(f"SSE: no se pudo re-procesar sesion: {exc}")

        if not chapters_to_process:
            chapters_to_process = chapters_from_request

        if not chapters_to_process:
            yield f'data: {json_module.dumps({"type": "error", "message": "No se pudo obtener el contenido de los capítulos"})}\n\n'
            return

        # Imágenes del documento (opt-in): se extraen una sola vez para todo
        # el documento y se reparten por capítulo según sus páginas. Solo se
        # adjuntan al primer chunk de cada capítulo (no a todos) para no
        # repetir la misma imagen en cada llamada — cada imagen mandada de
        # nuevo suma costo/latencia en el proveedor de IA.
        doc_images = []
        if include_images and session_file and os.path.exists(session_file):
            try:
                doc_images = extract_page_images(session_file, max_images=12)
            except Exception as exc:
                logger.warning(f"No se pudieron extraer imágenes del documento: {exc}")

        selected_tokens = _chapters_total_tokens(chapters_to_process)
        run_budget = settings.CONTENIDO_MAX_RUN_TOKENS
        if selected_tokens > run_budget:
            budget_msg = (
                f'Se seleccionaron {_fmt_es(selected_tokens)} tokens de contenido, y el máximo por '
                f'tanda es {_fmt_es(run_budget)}. Elegir menos capítulos y generar el resto en otra tanda.'
            )
            yield f'data: {json_module.dumps({"type": "error", "message": budget_msg})}\n\n'
            return

        # Pre-calcular total de chunks para progress
        chapter_splits = []
        total_chunks_all = 0
        for chapter in chapters_to_process:
            content = chapter.get('content', chapter.get('content_preview', ''))
            chunks = _split_into_chunks(content, max_tokens=content_chunk_tokens)
            chapter_splits.append(chunks)
            total_chunks_all += len(chunks)

        yield f'data: {json_module.dumps({"type": "start", "total_chunks": total_chunks_all, "filename": filename, "existing_count": len(existing_questions_list)})}\n\n'

        seen_keys = set(existing_texts_set)  # inicializar con preguntas ya en BD
        total_generated = 0
        chunk_idx_global = 0

        target_reached = False
        for chapter, chunks in zip(chapters_to_process, chapter_splits):
            if target_reached:
                break
            title = chapter.get('title', 'Capítulo')
            pages = chapter.get('pages', [])
            chapter_pages = set(pages)
            chapter_images = [img for img in doc_images if img['page'] in chapter_pages][:3]
            # Física → impresa (ver document_processor._detect_printed_page_number
            # y el mismo comentario en generate_questions_from_chapters): si no
            # se detectó número impreso para una página, se cita la física.
            printed_map = dict(zip(pages, chapter.get('printed_pages', [])))
            # Si el usuario pidió una cantidad fija por bloque, respetarla;
            # si no, distribuir total_questions entre todos los chunks. El piso
            # es 1 (no 2): con muchos chunks chicos, un piso de 2 podía duplicar
            # ampliamente lo pedido (ej. 20 preguntas en un documento con 30
            # chunks resultaba en 60, no 20).
            if questions_per_block_override > 0:
                questions_per_chunk = questions_per_block_override
            else:
                questions_per_chunk = max(1, min(12, total_questions // max(total_chunks_all, 1)))

            for i, chunk in enumerate(chunks):
                chunk_idx_global += 1
                chunk_images = chapter_images if (chapter_images and i == 0) else []
                chunk_has_text = _chunk_has_content(chunk['text'])
                if not chunk_has_text and not chunk_images:
                    logger.warning(f"SSE chunk {chunk_idx_global} de '{title}' sin texto suficiente, se omite.")
                    yield f'data: {json_module.dumps({"type": "chunk_error", "chunk": chunk_idx_global, "total_chunks": total_chunks_all, "chapter_title": title, "message": "Fragmento sin texto extraíble (posible página escaneada o solo con imágenes) — no se generaron preguntas."})}\n\n'
                    continue
                # Pequeño respiro entre requests: varios fragmentos seguidos sin
                # pausa pueden superar el límite por minuto (TPM/RPM) de Groq.
                if chunk_idx_global > 1:
                    time.sleep(2)
                try:
                    # Fragmento sin texto propio (página escaneada) pero con
                    # imágenes asociadas: no tiene sentido una llamada de solo
                    # texto — el chunk_error de arriba ya la habría descartado
                    # si tampoco tuviera imágenes.
                    text_questions = []
                    if chunk_has_text:
                        text_questions = _generate_questions_for_chunk(
                            chunk['text'], title, questions_per_chunk, i, len(chunks),
                            question_types=question_types, backend=backend,
                            existing_questions=existing_questions_list,
                            output_tokens_ceiling=output_tokens_ceiling,
                        )
                    image_questions = []
                    if chunk_images:
                        # Texto e imágenes van por separado a su proveedor
                        # correspondiente (ver DemoRoutingBackend en
                        # ai_router.py): Groq no tiene modelos con visión, así
                        # que mandar todo junto en una sola llamada forzaba
                        # también el texto a Gemini. Si esta segunda llamada
                        # falla, no se pierde lo ya generado del texto — se
                        # loguea y se sigue solo con eso.
                        try:
                            image_questions = _generate_questions_for_chunk(
                                chunk['text'], title, questions_per_chunk, i, len(chunks),
                                question_types=question_types, backend=backend,
                                existing_questions=existing_questions_list + text_questions,
                                images=chunk_images,
                                output_tokens_ceiling=output_tokens_ceiling,
                            )
                        except Exception as img_exc:
                            logger.warning(f"SSE chunk {chunk_idx_global} (imágenes) de '{title}' falló: {img_exc}")
                    raw_questions = text_questions + image_questions
                    # Páginas del FRAGMENTO puntual del que salieron estas
                    # preguntas (no de todo el capítulo/tanda) — ver comentario
                    # análogo en generate_questions_from_chapters.
                    chunk_pages = chunk['pages'] or pages
                    display_pages = sorted({printed_map.get(p) or p for p in chunk_pages}) if chunk_pages else []
                    new_qs = []
                    for q in raw_questions:
                        key = q.get('pregunta', '').lower().strip()[:80]
                        if key and key not in seen_keys:
                            seen_keys.add(key)
                            q['source_chapters'] = [{'title': title, 'pages': display_pages}]
                            q['source_file'] = filename
                            new_qs.append(q)

                    # Tope duro: "cantidad de preguntas" (total_questions) es un
                    # techo absoluto, tanto si el modelo no respeta "generá
                    # exactamente N preguntas" al pie de la letra, como si el
                    # usuario configuró "preguntas por bloque" — ese campo solo
                    # controla el tamaño de cada bloque que se muestra en pantalla
                    # (útil en modo pausa), nunca debería poder superar el total
                    # pedido. Antes, con un override de bloque > 0, el total se
                    # ignoraba por completo.
                    remaining = max(0, total_questions - total_generated)
                    new_qs = new_qs[:remaining]

                    total_generated += len(new_qs)
                    event = {
                        'type': 'questions',
                        'questions': new_qs,
                        'chunk': chunk_idx_global,
                        'total_chunks': total_chunks_all,
                        'chapter_title': title,
                    }
                    yield f'data: {json_module.dumps(event)}\n\n'

                    # "cantidad de preguntas" es un objetivo total, no un piso:
                    # paramos apenas lo alcanzamos en vez de seguir procesando el
                    # resto del documento y generar de más.
                    if total_generated >= total_questions:
                        target_reached = True
                        break

                except GeneratorExit:
                    return
                except Exception as exc:
                    logger.warning(f"SSE chunk {chunk_idx_global} error: {exc}")
                    yield f'data: {json_module.dumps({"type": "chunk_error", "chunk": chunk_idx_global, "total_chunks": total_chunks_all, "chapter_title": title, "message": str(exc)})}\n\n'

        yield f'data: {json_module.dumps({"type": "done", "total": total_generated})}\n\n'

    with _jobs_lock:
        job = _jobs.pop(job_id, None)

    if not job:
        def _not_found():
            import json as j
            yield f'data: {j.dumps({"type": "error", "message": "Job no encontrado o expirado"})}\n\n'
        response = StreamingHttpResponse(_not_found(), content_type='text/event-stream')
        response['Cache-Control'] = 'no-cache'
        return response

    from .ai_router import get_backend_for_user
    _sse_backend = get_backend_for_user(request.user)

    response = StreamingHttpResponse(
        event_stream(job, request.user.id, _sse_backend),
        content_type='text/event-stream'
    )
    response['Cache-Control'] = 'no-cache'
    response['X-Accel-Buffering'] = 'no'
    return response


@login_required
@require_http_methods(["GET"])
def document_page_preview(request):
    """
    Devuelve metadata del documento en sesión para el visor de páginas:
    - Para PDF: número total de páginas y URL para servirlo en PDF.js
    - Para DOCX: bloques de texto por sección (para docx-preview + checkboxes)
    - Para PPTX: texto de cada slide (tarjeta por slide con checkbox)

    GET params: ninguno (usa la sesión doc_processor)

    Returns JSON:
        {
          "success": true,
          "file_type": "pdf"|"docx"|"pptx"|"txt",
          "filename": "...",
          "total_pages": N,          # PDF
          "file_url": "/media/...",  # PDF — para PDF.js
          "slides": [...],           # PPTX
          "sections": [...],         # DOCX/TXT
        }
    """
    doc_session = request.session.get('doc_processor', {})
    file_path = doc_session.get('file_path', '')
    filename = doc_session.get('filename', '')

    if not file_path or not os.path.exists(file_path):
        # Fallback: intentar recuperar desde contenido_id si la sesión fue
        # pisada por una condición de carrera (SESSION_SAVE_EVERY_REQUEST).
        contenido_id = request.GET.get('contenido_id', '').strip()
        if contenido_id:
            try:
                from material.models import Contenido
                contenido = Contenido.objects.get(id=int(contenido_id), uploaded_by=request.user)
                candidate_path = contenido.file.path if contenido.file else ''
                if candidate_path and os.path.exists(candidate_path):
                    file_path = candidate_path
                    filename = os.path.basename(file_path)
                    # Re-establecer la sesión para llamadas posteriores
                    request.session.setdefault('doc_processor', {})
                    request.session['doc_processor']['file_path'] = file_path
                    request.session['doc_processor']['filename'] = filename
                    request.session.modified = True
            except (Contenido.DoesNotExist, ValueError, AttributeError, OSError):
                pass

        if not file_path or not os.path.exists(file_path):
            return JsonResponse({'success': False, 'error': 'No hay documento en sesión. Sube o selecciona uno primero.'}, status=400)

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == '.pdf':
            doc = fitz.open(file_path)
            total_pages = doc.page_count
            doc.close()
            # Usar endpoint interno para servir el archivo (funciona con DEBUG=False en Render)
            file_url = '/doc-processor/serve-file/'
            return JsonResponse({
                'success': True,
                'file_type': 'pdf',
                'filename': filename,
                'total_pages': total_pages,
                'file_url': file_url,
            })

        elif ext == '.pptx':
            from pptx import Presentation as _Prs
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            import base64
            prs = _Prs(file_path)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text.strip())
                text = '\n'.join(texts)
                # Sin texto: buscar la imagen embebida más grande del slide
                # (puede haber varias, ej. íconos sueltos) para mostrar algo
                # visual en el preview en vez del hueco vacío. No cubre
                # diagramas/formas nativas (SmartArt, autoshapes) — eso
                # requeriría renderizar el slide completo, que no tenemos
                # forma de hacer sin instalar LibreOffice en Render.
                image_data_uri = None
                if not text:
                    try:
                        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
                        if pictures:
                            largest = max(pictures, key=lambda s: (s.width or 0) * (s.height or 0))
                            image = largest.image
                            image_data_uri = f'data:{image.content_type};base64,{base64.b64encode(image.blob).decode("ascii")}'
                    except Exception:
                        image_data_uri = None
                slides.append({
                    'slide_number': i,
                    'text': text or (f'(Slide {i} sin texto)' if not image_data_uri else ''),
                    'char_count': sum(len(t) for t in texts),
                    'image': image_data_uri,
                })
            return JsonResponse({
                'success': True,
                'file_type': 'pptx',
                'filename': filename,
                'total_pages': len(slides),
                'slides': slides,
            })

        elif ext == '.docx':
            from docx import Document as _Doc
            doc = _Doc(file_path)

            # ── OPCIÓN D: páginas virtuales por cantidad de caracteres ──────────
            # ROLLBACK: reemplazar este bloque con el bloque comentado de abajo
            # para volver al modo de secciones por heading.
            DOCX_PAGE_SIZE = 2800  # chars ≈ 1 página A4 en texto normal
            sections = []
            current_text = []
            current_chars = 0
            page_idx = 1
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                current_text.append(text)
                current_chars += len(text) + 1
                if current_chars >= DOCX_PAGE_SIZE:
                    sections.append({
                        'section_number': page_idx,
                        'title': f'Página {page_idx}',
                        'text': '\n'.join(current_text),
                        'char_count': current_chars,
                    })
                    page_idx += 1
                    current_text = []
                    current_chars = 0
            if current_text:
                sections.append({
                    'section_number': page_idx,
                    'title': f'Página {page_idx}',
                    'text': '\n'.join(current_text),
                    'char_count': current_chars,
                })
            # ── FIN OPCIÓN D ─────────────────────────────────────────────────────

            # [ROLLBACK DOCX SECTIONS — secciones por heading, descommentar para revertir]
            # sections = []
            # current_text = []
            # current_chars = 0
            # section_idx = 1
            # heading_title = 'Inicio'
            # for para in doc.paragraphs:
            #     text = para.text.strip()
            #     if not text:
            #         continue
            #     if para.style.name.startswith('Heading'):
            #         if current_text:
            #             sections.append({'section_number': section_idx, 'title': heading_title,
            #                              'text': '\n'.join(current_text), 'char_count': current_chars})
            #             section_idx += 1; current_text = []; current_chars = 0
            #         heading_title = text
            #     else:
            #         current_text.append(text); current_chars += len(text)
            #         if current_chars > 1200:
            #             sections.append({'section_number': section_idx, 'title': heading_title,
            #                              'text': '\n'.join(current_text), 'char_count': current_chars})
            #             section_idx += 1; current_text = []; current_chars = 0
            # if current_text:
            #     sections.append({'section_number': section_idx, 'title': heading_title,
            #                      'text': '\n'.join(current_text), 'char_count': current_chars})
            # URL del archivo para docx-preview.js
            file_url = '/doc-processor/serve-file/'
            return JsonResponse({
                'success': True,
                'file_type': 'docx',
                'filename': filename,
                'total_pages': len(sections),
                'file_url': file_url,
                'sections': sections,
            })

        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                lines = f.readlines()
            # Dividir en bloques de ~50 líneas
            block_size = 50
            sections = []
            for i in range(0, len(lines), block_size):
                block = ''.join(lines[i:i + block_size]).strip()
                sections.append({
                    'section_number': i // block_size + 1,
                    'title': f'Líneas {i+1}–{min(i+block_size, len(lines))}',
                    'text': block,
                    'char_count': len(block),
                })
            return JsonResponse({
                'success': True,
                'file_type': 'txt',
                'filename': filename,
                'total_pages': len(sections),
                'sections': sections,
            })

        else:
            return JsonResponse({'success': False, 'error': f'Formato no soportado para preview: {ext}'}, status=400)

    except Exception as e:
        logger.exception("Error en document_page_preview")
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


@login_required
@require_http_methods(["POST"])
def get_pages_text(request):
    """
    Dado un conjunto de números de página/slide/sección seleccionados por el usuario,
    extrae el texto de esas unidades del documento en sesión.
    Devuelve el mismo formato de "chapters" que usa generate_questions_from_chapters.

    POST JSON:
        {
          "pages": [1, 3, 5],        # para PDF: números de página (1-based)
          "slides": [2, 4],          # para PPTX: números de slide (1-based)
          "sections": [1, 2],        # para DOCX/TXT: números de sección (1-based)
        }
    Returns JSON:
        {
          "success": true,
          "chapters": [{"title": "...", "content": "...", "tokens": N, "pages": [...]}]
        }
    """
    import json as json_module
    try:
        data = json_module.loads(request.body)
    except Exception:
        return JsonResponse({'success': False, 'error': 'JSON inválido'}, status=400)

    doc_session = request.session.get('doc_processor', {})
    file_path = doc_session.get('file_path', '')
    filename = doc_session.get('filename', '')

    if not file_path or not os.path.exists(file_path):
        contenido_id = data.get('contenido_id', '')
        if contenido_id:
            try:
                from material.models import Contenido
                contenido = Contenido.objects.get(id=int(contenido_id), uploaded_by=request.user)
                candidate_path = contenido.file.path if contenido.file else ''
                if candidate_path and os.path.exists(candidate_path):
                    file_path = candidate_path
                    filename = os.path.basename(file_path)
                    request.session.setdefault('doc_processor', {})
                    request.session['doc_processor']['file_path'] = file_path
                    request.session['doc_processor']['filename'] = filename
                    request.session.modified = True
            except (Contenido.DoesNotExist, ValueError, AttributeError, OSError):
                pass

        if not file_path or not os.path.exists(file_path):
            return JsonResponse({'success': False, 'error': 'No hay documento en sesión.'}, status=400)

    ext = os.path.splitext(file_path)[1].lower()

    try:
        chapters = []

        if ext == '.pdf':
            selected_pages = [int(p) for p in data.get('pages', [])]
            if not selected_pages:
                return JsonResponse({'success': False, 'error': 'No se enviaron páginas.'}, status=400)
            doc = fitz.open(file_path)
            for page_num in sorted(selected_pages):
                if 1 <= page_num <= doc.page_count:
                    page = doc[page_num - 1]
                    text = page.get_text().strip()
                    if text:
                        chapters.append({
                            'title': f'Página {page_num}',
                            'content': text,
                            'tokens': count_tokens(text),
                            'pages': [page_num],
                        })
            doc.close()

        elif ext == '.pptx':
            from pptx import Presentation as _Prs
            selected = set(int(s) for s in data.get('slides', []))
            if not selected:
                return JsonResponse({'success': False, 'error': 'No se enviaron slides.'}, status=400)
            prs = _Prs(file_path)
            for i, slide in enumerate(prs.slides, 1):
                if i in selected:
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            texts.append(shape.text.strip())
                    text = '\n'.join(texts)
                    if text:
                        chapters.append({
                            'title': f'Slide {i}',
                            'content': text,
                            'tokens': count_tokens(text),
                            'pages': [i],
                        })

        elif ext in ('.docx', '.txt'):
            selected = set(int(s) for s in data.get('sections', []))
            if not selected:
                return JsonResponse({'success': False, 'error': 'No se enviaron secciones.'}, status=400)
            # Re-usar la misma lógica de document_page_preview para extraer secciones
            if ext == '.docx':
                from docx import Document as _Doc
                doc = _Doc(file_path)

                # ── OPCIÓN D: páginas virtuales (debe coincidir con document_page_preview) ──
                # ROLLBACK: reemplazar con el bloque comentado al final de este if
                DOCX_PAGE_SIZE = 2800
                sections_all = []
                current_text = []
                current_chars = 0
                page_idx = 1
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    current_text.append(text)
                    current_chars += len(text) + 1
                    if current_chars >= DOCX_PAGE_SIZE:
                        sections_all.append((page_idx, f'Página {page_idx}', '\n'.join(current_text)))
                        page_idx += 1
                        current_text = []
                        current_chars = 0
                if current_text:
                    sections_all.append((page_idx, f'Página {page_idx}', '\n'.join(current_text)))
                # ── FIN OPCIÓN D ────────────────────────────────────────────────────

                # [ROLLBACK DOCX SECTIONS — descommentar para revertir]
                # sections_all = []
                # current_text = []; current_chars = 0; section_idx = 1; heading_title = 'Inicio'
                # for para in doc.paragraphs:
                #     text = para.text.strip()
                #     if not text: continue
                #     if para.style.name.startswith('Heading'):
                #         if current_text:
                #             sections_all.append((section_idx, heading_title, '\n'.join(current_text)))
                #             section_idx += 1; current_text = []; current_chars = 0
                #         heading_title = text
                #     else:
                #         current_text.append(text); current_chars += len(text)
                #         if current_chars > 1200:
                #             sections_all.append((section_idx, heading_title, '\n'.join(current_text)))
                #             section_idx += 1; current_text = []; current_chars = 0
                # if current_text:
                #     sections_all.append((section_idx, heading_title, '\n'.join(current_text)))
            else:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    lines = f.readlines()
                block_size = 50
                sections_all = []
                for i in range(0, len(lines), block_size):
                    block = ''.join(lines[i:i + block_size]).strip()
                    idx = i // block_size + 1
                    sections_all.append((idx, f'Líneas {i+1}–{min(i+block_size, len(lines))}', block))

            for sec_num, title, text in sections_all:
                if sec_num in selected and text:
                    chapters.append({
                        'title': title,
                        'content': text,
                        'tokens': count_tokens(text),
                        'pages': [sec_num],
                    })

        if not chapters:
            return JsonResponse({'success': False, 'error': 'No se pudo extraer texto de las unidades seleccionadas.'}, status=400)

        return JsonResponse({
            'success': True,
            'filename': filename,
            'chapters': chapters,
            'total_tokens': sum(ch['tokens'] for ch in chapters),
        })

    except Exception as e:
        logger.exception("Error en get_pages_text")
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


@login_required
@require_http_methods(["POST"])
def save_generated_questions(request):
    """
    Guarda preguntas generadas por IA en la base de datos.
    
    POST JSON:
        - approved: lista de preguntas aprobadas
        - rejected: lista de preguntas rechazadas
        - filename: nombre del archivo fuente
    
    Returns:
        JSON con resultado de la operación
    """
    import json as json_module
    from material.models import Question, Subject, Topic, Contenido
    
    try:
        data = json_module.loads(request.body)
        approved = data.get('approved', [])
        rejected = data.get('rejected', [])
        filename = data.get('filename', 'Documento')
        subject_ids = data.get('subject_ids', [])
        contenido_id = data.get('contenido_id')

        # Resolver el Contenido de origen si viene informado
        contenido_origen = None
        if contenido_id:
            try:
                contenido_origen = Contenido.objects.get(id=contenido_id, uploaded_by=request.user)
            except Contenido.DoesNotExist:
                pass
        
        saved_count = 0
        topic_id = data.get('topic_id')
        subtopic_id = data.get('subtopic_id')
        new_topic_name = (data.get('new_topic_name') or '').strip()
        new_subtopic_name = (data.get('new_subtopic_name') or '').strip()

        # Resolver materias seleccionadas por el usuario
        selected_subjects = list(Subject.objects.filter(id__in=subject_ids)) if subject_ids else []
        if not selected_subjects:
            # Fallback: la materia ya asociada al Contenido de origen (la que
            # el usuario eligió al subir el documento), NO la primera materia
            # del sistema en orden alfabético — eso hacía que, con contenido
            # semilla cargado, cualquier guardado sin materia explícita
            # terminara clasificado en la materia que alfabéticamente
            # apareciera primero (p. ej. "Bases de Datos"), sin relación con
            # lo que el usuario estaba trabajando.
            fallback = contenido_origen.subjects.first() if contenido_origen else None
            if not fallback:
                # Excluye materias semilla (is_seed_demo) del fallback por la
                # misma razón que el comentario de arriba: sin esto, contenido
                # real sin materia explícita podía terminar clasificado en la
                # materia de ejemplo del asistente (primera alfabéticamente).
                fallback = Subject.objects.filter(is_seed_demo=False).first()
            if not fallback:
                return JsonResponse({
                    'success': False,
                    'error': 'No hay materias configuradas en el sistema'
                }, status=400)
            selected_subjects = [fallback]

        # El tema/subtema se ancla a la primera materia seleccionada
        from material.models import Subtopic
        default_subject = selected_subjects[0]

        # Resolver tema
        default_topic = None
        if topic_id:
            try:
                default_topic = Topic.objects.get(id=topic_id, subject=default_subject)
            except Topic.DoesNotExist:
                pass
        if not default_topic and new_topic_name:
            default_topic, _ = Topic.objects.get_or_create(
                name=new_topic_name,
                subject=default_subject
            )
        if not default_topic:
            # Fallback automático
            default_topic, _ = Topic.objects.get_or_create(
                name=f"Preguntas de {filename}",
                subject=default_subject
            )

        # Resolver subtema
        default_subtopic = None
        if subtopic_id:
            try:
                default_subtopic = Subtopic.objects.get(id=subtopic_id, topic=default_topic)
            except Subtopic.DoesNotExist:
                pass
        if not default_subtopic and new_subtopic_name:
            default_subtopic, _ = Subtopic.objects.get_or_create(
                name=new_subtopic_name,
                topic=default_topic
            )

        # Idempotencia: un reintento de red o un doble envío (ver Fix B en el
        # frontend, botón de generar) puede mandar este POST más de una vez
        # con las mismas preguntas — Question no tiene unique_together, así
        # que sin esto se duplicaban filas. Se compara texto normalizado
        # dentro del mismo Contenido de origen (o la(s) materia(s) elegidas,
        # si no hay Contenido), no de forma global.
        existing_questions_qs = Question.objects.filter(user=request.user, generated_by_ai=True)
        if contenido_origen:
            existing_questions_qs = existing_questions_qs.filter(contenido=contenido_origen)
        else:
            existing_questions_qs = existing_questions_qs.filter(contenido__isnull=True, subjects__in=selected_subjects)
        existing_question_texts = set(
            t.strip() for t in existing_questions_qs.values_list('question_text', flat=True)
        )
        skipped_duplicates = 0

        # Guardar preguntas aprobadas
        for q_data in approved:
            question_text = (q_data.get('pregunta') or '').strip()
            if question_text and question_text in existing_question_texts:
                skipped_duplicates += 1
                continue
            question = Question(
                topic=default_topic,
                subtopic=default_subtopic,
                question_type=q_data.get('tipo', 'opcion_multiple'),
                question_text=q_data.get('pregunta', ''),
                answer_text=q_data.get('respuesta', ''),
                difficulty=q_data.get('dificultad', 3),
                bloom_level=q_data.get('bloom_nivel') or None,
                user=request.user,
                generated_by_ai=True,
                ai_approved=True,
                contenido=contenido_origen
            )

            # Guardar opciones si existen
            if 'opciones' in q_data:
                question.options = q_data['opciones']

            # Guardar información de capítulos fuente
            if 'source_chapters' in q_data:
                question.source_chapters = q_data['source_chapters']
                question.source_page = _first_source_page(q_data['source_chapters'])

            question.save()
            question.subjects.set(selected_subjects)
            saved_count += 1
            if question_text:
                existing_question_texts.add(question_text)

        # Guardar preguntas rechazadas (para registro)
        for q_data in rejected:
            question_text = (q_data.get('pregunta') or '').strip()
            if question_text and question_text in existing_question_texts:
                skipped_duplicates += 1
                continue
            question = Question(
                topic=default_topic,
                subtopic=default_subtopic,
                question_type=q_data.get('tipo', 'opcion_multiple'),
                question_text=q_data.get('pregunta', ''),
                answer_text=q_data.get('respuesta', ''),
                difficulty=q_data.get('dificultad', 3),
                bloom_level=q_data.get('bloom_nivel') or None,
                user=request.user,
                generated_by_ai=True,
                ai_approved=False,
                contenido=contenido_origen
            )

            if 'opciones' in q_data:
                question.options = q_data['opciones']

            if 'source_chapters' in q_data:
                question.source_chapters = q_data['source_chapters']
                question.source_page = _first_source_page(q_data['source_chapters'])

            question.save()
            question.subjects.set(selected_subjects)
            if question_text:
                existing_question_texts.add(question_text)

        return JsonResponse({
            'success': True,
            'saved_count': saved_count,
            'approved_count': len(approved),
            'rejected_count': len(rejected),
            'total_count': len(approved) + len(rejected),
            'skipped_duplicates': skipped_duplicates
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e)
        }, status=500)


@login_required
@require_http_methods(["GET"])
def get_topics_by_subject(request, subject_id):
    """Retorna temas y subtemas para una materia — usado por el modal de guardado."""
    from material.models import Topic
    try:
        topics = (
            Topic.objects
            .filter(subject_id=subject_id)
            .order_by('name')
            .prefetch_related('subtopic_set')
        )
        result = [
            {
                'id': t.id,
                'name': t.name,
                'subtopics': [
                    {'id': s.id, 'name': s.name}
                    for s in t.subtopic_set.all().order_by('name')
                ],
            }
            for t in topics
        ]
        return JsonResponse({'success': True, 'topics': result})
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)


@login_required
def serve_doc_file(request):
    """
    Sirve el archivo del documento en sesión directamente via Django.
    Necesario en producción (Render) donde DEBUG=False y /media/ no es servido.
    Solo accesible para el usuario dueño de la sesión.
    """
    doc_session = request.session.get('doc_processor', {})
    file_path = doc_session.get('file_path', '')

    if not file_path or not os.path.exists(file_path):
        raise Http404("Archivo no disponible")

    # Validar que el path esté dentro de MEDIA_ROOT (seguridad)
    try:
        real_path = os.path.realpath(file_path)
        media_real = os.path.realpath(settings.MEDIA_ROOT)
        if not real_path.startswith(media_real):
            raise Http404("Acceso denegado")
    except Exception:
        raise Http404("Archivo no disponible")

    ext = os.path.splitext(file_path)[1].lower()
    content_types = {
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.txt': 'text/plain',
    }
    content_type = content_types.get(ext, 'application/octet-stream')
    response = FileResponse(open(file_path, 'rb'), content_type=content_type)
    response['Content-Disposition'] = f'inline; filename="{os.path.basename(file_path)}"'
    return response


