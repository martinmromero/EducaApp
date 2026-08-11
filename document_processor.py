"""
Document Processor Module
==========================
Módulo para procesar documentos (PDF, DOCX, PPTX, TXT) con foco en:
- Extracción de estructura jerárquica (capítulos, secciones)
- Limpieza de texto repetitivo (headers, footers, números de página)
- Optimización de tokens para envío a IA (GPT-4, Claude, etc.)

Autor: EducaApp
Fecha: 29 octubre 2025
"""

import fitz  # PyMuPDF
import tiktoken
from markdownify import markdownify as md
from docx import Document
from pptx import Presentation
from typing import Dict, List, Tuple, Optional
import re
from collections import Counter
import json


class DocumentProcessor:
    """
    Procesador universal de documentos con detección de estructura
    y optimización para envío a IA.
    """
    
    def __init__(self, encoding_name: str = "cl100k_base"):
        """
        Inicializa el procesador.
        
        Args:
            encoding_name: Nombre del encoding de tiktoken (ej: "cl100k_base" para GPT-4)
        """
        self.encoding = tiktoken.get_encoding(encoding_name)
        self.stats = {
            'total_pages': 0,
            'total_tokens': 0,
            'chapters': [],
            'removed_headers': 0,
            'removed_footers': 0
        }
    
    # ========================================================================
    # PDF PROCESSING (PyMuPDF)
    # ========================================================================
    
    def process_pdf(self, file_path: str,
                   remove_headers: bool = True,
                   remove_footers: bool = True,
                   extract_toc: bool = True,
                   max_pages: Optional[int] = None,
                   pages_per_block: int = 40) -> Dict:
        """
        Procesa un PDF completo y extrae estructura + contenido limpio.

        Args:
            file_path: Ruta al archivo PDF
            remove_headers: Si True, detecta y elimina headers repetitivos
            remove_footers: Si True, detecta y elimina footers repetitivos
            extract_toc: Si True, extrae tabla de contenidos si existe
            max_pages: Techo de sanity (no de "tamaño de libro"): rechaza
                documentos con más páginas antes de extraer texto. Extraer
                texto es liviano incluso en libros largos, así que este
                valor solo protege contra PDFs patológicos/corruptos, no
                limita subir un libro completo.
            pages_per_block: Cuando el PDF no tiene tabla de contenidos
                (TOC), en vez de tratar todo el documento como un único
                "capítulo" gigante (imposible de seleccionar por partes),
                se lo divide en bloques de esta cantidad de páginas.

        Returns:
            Diccionario con estructura:
            {
                'metadata': {...},
                'toc': [...],
                'chapters': [{'title': str, 'content': str, 'tokens': int, 'pages': [int]}],
                'stats': {...}
            }
        """
        result = {
            'metadata': {},
            'toc': [],
            'chapters': [],
            'full_text': '',
            'stats': {}
        }

        # Abrir PDF con PyMuPDF
        doc = fitz.open(file_path)

        if max_pages and doc.page_count > max_pages:
            page_count = doc.page_count
            doc.close()
            raise ValueError(
                f'El documento tiene {page_count} páginas y supera el máximo permitido '
                f'de {max_pages}. Subí solo las páginas o el capítulo que necesitás analizar.'
            )

        # Extraer metadata
        result['metadata'] = {
            'title': doc.metadata.get('title', ''),
            'author': doc.metadata.get('author', ''),
            'subject': doc.metadata.get('subject', ''),
            'total_pages': doc.page_count,
            'format': doc.metadata.get('format', 'PDF')
        }
        
        self.stats['total_pages'] = doc.page_count
        
        # Extraer TOC (tabla de contenidos) si existe
        if extract_toc:
            toc = doc.get_toc()
            result['toc'] = [
                {'level': level, 'title': title, 'page': page}
                for level, title, page in toc
            ]
        
        # Detectar headers/footers repetitivos
        headers_to_remove = []
        footers_to_remove = []

        if remove_headers or remove_footers:
            headers_to_remove, footers_to_remove = self._detect_repetitive_text(doc)
        
        # Extraer texto por capítulo si hay TOC, sino por páginas
        if result['toc'] and len(result['toc']) > 0:
            result['chapters'] = self._extract_chapters_from_toc(
                doc, result['toc'], headers_to_remove, footers_to_remove
            )
        elif doc.page_count > pages_per_block:
            # Sin TOC y documento largo: partirlo en bloques de páginas para
            # que se pueda seleccionar y generar por partes, en vez de un
            # único capítulo gigante imposible de procesar de una corrida.
            result['chapters'] = self._extract_chapters_by_page_blocks(
                doc, headers_to_remove, footers_to_remove, pages_per_block
            )
        else:
            # Sin TOC y documento corto: un solo capítulo alcanza. Se
            # extrae directo por página (una sola pasada sobre el PDF)
            # para poder incrustar los marcadores de página.
            doc_pages = []
            doc_printed_pages = []
            page_texts = []
            for page_num in range(doc.page_count):
                raw_text = doc[page_num].get_text()
                doc_pages.append(page_num + 1)
                doc_printed_pages.append(self._detect_printed_page_number(raw_text))
                text = self._remove_repetitive_patterns(raw_text, headers_to_remove, footers_to_remove)
                page_texts.append(text.strip())
            content = self._join_pages_with_markers(doc_pages, page_texts)

            result['chapters'] = [{
                'title': 'Documento completo',
                'content': content,
                'tokens': self.count_tokens(content),
                'pages': doc_pages,
                'printed_pages': doc_printed_pages,
            }]
        
        # Calcular stats finales
        result['stats'] = {
            'total_pages': doc.page_count,
            'total_chapters': len(result['chapters']),
            'total_tokens': sum(ch['tokens'] for ch in result['chapters']),
            'removed_headers': self.stats['removed_headers'],
            'removed_footers': self.stats['removed_footers']
        }
        
        doc.close()
        return result
    
    def _detect_repetitive_text(self, doc: fitz.Document,
                                threshold: float = 0.7) -> Tuple[List[str], List[str]]:
        """
        Detecta headers y footers repetitivos.

        Reusa el `fitz.Document` ya abierto por process_pdf en vez de abrir
        el PDF de nuevo con pdfplumber: pdfplumber (pdfminer.six por debajo,
        con análisis de layout completo) es notablemente más lento que
        PyMuPDF para lo mismo, y aquí solo se necesita texto plano de las
        primeras páginas — no vale la pena una segunda librería ni una
        segunda apertura del archivo solo para eso.

        Args:
            doc: Documento PyMuPDF ya abierto
            threshold: Porcentaje de páginas donde debe aparecer para considerarse repetitivo (0-1)

        Returns:
            Tupla (headers, footers) con listas de textos a eliminar
        """
        headers = []
        footers = []

        total_pages = doc.page_count
        if total_pages == 0:
            return [], []

        # Recolectar primeras/últimas líneas de cada página
        top_lines = []
        bottom_lines = []

        for page_num in range(min(10, total_pages)):  # Samplear primeras 10 páginas
            text = doc[page_num].get_text()
            if not text:
                continue

            lines = text.strip().split('\n')
            if len(lines) > 0:
                top_lines.append(lines[0].strip())
            if len(lines) > 1:
                bottom_lines.append(lines[-1].strip())

        # Contar frecuencias
        top_counter = Counter(top_lines)
        bottom_counter = Counter(bottom_lines)

        # Identificar repetitivos (aparecen en >threshold% de páginas)
        min_occurrences = int(len(top_lines) * threshold)

        for text, count in top_counter.items():
            if count >= min_occurrences and len(text) > 5:  # Mín 5 chars
                headers.append(text)
                self.stats['removed_headers'] += count

        for text, count in bottom_counter.items():
            if count >= min_occurrences and len(text) > 3:
                # Ignorar números de página simples
                if not re.match(r'^\d+$', text):
                    footers.append(text)
                    self.stats['removed_footers'] += count

        return headers, footers
    
    def _extract_chapters_from_toc(self, doc: fitz.Document,
                                   toc: List[Dict],
                                   headers: List[str],
                                   footers: List[str]) -> List[Dict]:
        """
        Extrae contenido organizado por capítulos según TOC.
        Detecta automáticamente el nivel más granular: si un nivel-1 tiene
        hijos en nivel-2, usa esos hijos como capítulos (con prefijo de la
        parte). Si un nivel-1 no tiene hijos, lo considera capítulo directo.
        """
        chapters = []

        flat_items = self._flatten_toc_to_chapters(toc)

        # Fallback: si no se obtuvo nada, usar nivel-1 original
        if not flat_items:
            flat_items = [
                {'display_title': item['title'].strip(), 'page': item['page']}
                for item in toc if item['level'] == 1
            ]

        for i, item in enumerate(flat_items):
            start_page = max(0, item['page'] - 1)  # PyMuPDF índice 0

            # Página final: inicio del siguiente ítem o fin del documento
            if i < len(flat_items) - 1:
                end_page = flat_items[i + 1]['page'] - 1
            else:
                end_page = doc.page_count

            end_page = max(start_page + 1, min(end_page, doc.page_count))

            # Extraer texto de páginas del capítulo
            chapter_text = []
            chapter_pages = []
            chapter_printed_pages = []
            for page_num in range(start_page, end_page):
                if page_num < doc.page_count:
                    raw_text = doc[page_num].get_text()
                    chapter_pages.append(page_num + 1)
                    chapter_printed_pages.append(self._detect_printed_page_number(raw_text))
                    text = self._remove_repetitive_patterns(raw_text, headers, footers)
                    chapter_text.append(text.strip())

            # Sin texto (páginas escaneadas sin OCR, secciones vacías): no
            # tiene sentido ofrecerlo como capítulo seleccionable — generar
            # preguntas de ahí no tiene de dónde salir más que del propio
            # prompt. Ver [[project_fotosintesis_prompt_leak]]. Se chequea
            # ANTES de sumar los marcadores de página (ver _join_pages_with_
            # markers): esos marcadores no son texto real y siempre dejarían
            # el content "no vacío" aunque las páginas no tengan nada.
            if not ''.join(chapter_text).strip():
                continue

            content = self._join_pages_with_markers(chapter_pages, chapter_text)

            chapters.append({
                'title': item['display_title'],
                'content': content,
                'tokens': self.count_tokens(content),
                'pages': chapter_pages,
                'printed_pages': chapter_printed_pages,
            })

        return chapters

    def _extract_chapters_by_page_blocks(self, doc: fitz.Document,
                                         headers: List[str],
                                         footers: List[str],
                                         pages_per_block: int) -> List[Dict]:
        """
        Divide un PDF sin TOC en bloques consecutivos de páginas, cada uno
        tratado como un "capítulo" seleccionable. Es el equivalente sintético
        de _extract_chapters_from_toc para documentos sin estructura
        detectable (muy común en PDFs escaneados/convertidos).
        """
        chapters = []
        total_pages = doc.page_count

        for start_page in range(0, total_pages, pages_per_block):
            end_page = min(start_page + pages_per_block, total_pages)

            block_text = []
            block_pages = []
            block_printed_pages = []
            for page_num in range(start_page, end_page):
                raw_text = doc[page_num].get_text()
                block_pages.append(page_num + 1)
                block_printed_pages.append(self._detect_printed_page_number(raw_text))
                text = self._remove_repetitive_patterns(raw_text, headers, footers)
                block_text.append(text.strip())

            # Mismo criterio que _extract_chapters_from_toc: un bloque de
            # páginas enteramente escaneadas/sin texto no se ofrece como
            # "capítulo" seleccionable (chequeado antes de los marcadores de
            # página, que no son texto real — ver _join_pages_with_markers).
            if not ''.join(block_text).strip():
                continue

            content = self._join_pages_with_markers(block_pages, block_text)

            chapters.append({
                'title': f'Páginas {start_page + 1}-{end_page}',
                'content': content,
                'tokens': self.count_tokens(content),
                'pages': block_pages,
                'printed_pages': block_printed_pages,
            })

        return chapters

    def _flatten_toc_to_chapters(self, toc: List[Dict]) -> List[Dict]:
        """
        Aplana el TOC eligiendo el nivel más granular por sección.

        Reglas:
        - Si un nivel-1 tiene hijos en nivel-2 → reemplaza al nivel-1 con
          sus hijos de nivel-2, usando 'PARTE › Capítulo' como título.
          Los niveles 3+ dentro de un nivel-2 se ignoran como separadores
          (su contenido queda incluido en el rango de páginas del nivel-2).
        - Si un nivel-1 NO tiene hijos → se conserva como capítulo directo.
        - Entradas de nivel > 1 sin padre nivel-1 previo → se agregan tal cual.
        """
        flat_items = []
        i = 0

        while i < len(toc):
            item = toc[i]
            level = item['level']

            if level == 1:
                parent_title = item['title'].strip()
                i += 1  # avanzar al siguiente ítem

                children_added = 0
                # Consumir todos los descendientes (nivel > 1)
                while i < len(toc) and toc[i]['level'] > 1:
                    child = toc[i]
                    if child['level'] == 2:
                        flat_items.append({
                            'display_title': f"{parent_title} \u203a {child['title'].strip()}",
                            'page': child['page'],
                        })
                        children_added += 1
                    # nivel 3+: ignorar como separador (incluido en rango del nivel-2)
                    i += 1

                if children_added == 0:
                    # Nivel-1 sin hijos nivel-2 → usarlo como capítulo
                    flat_items.append({
                        'display_title': parent_title,
                        'page': item['page'],
                    })
            else:
                # Nivel > 1 sin padre nivel-1 previo → agregar directamente
                flat_items.append({
                    'display_title': item['title'].strip(),
                    'page': item['page'],
                })
                i += 1

        return flat_items

    
    def _remove_repetitive_patterns(self, text: str, 
                                    headers: List[str], 
                                    footers: List[str]) -> str:
        """
        Elimina patrones repetitivos del texto.
        """
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            line_stripped = line.strip()
            
            # Saltar headers
            if any(line_stripped.startswith(h) or line_stripped == h for h in headers):
                continue
            
            # Saltar footers
            if any(line_stripped.endswith(f) or line_stripped == f for f in footers):
                continue
            
            # Saltar números de página solitarios
            if re.match(r'^\d+$', line_stripped):
                continue
            
            cleaned_lines.append(line)

        return '\n'.join(cleaned_lines)

    # Patrones conservadores para el número "impreso" de una página (el que
    # el libro muestra en el pie/encabezado, ej. "322"), que puede no
    # coincidir con el índice físico de la página dentro del PDF subido (ej.
    # subir solo un capítulo aislado de un libro: la página física 2 del PDF
    # puede decir "322" impreso). Solo se buscan en la primera/última línea
    # no vacía de la página — donde vive un número de página real — para no
    # confundir un número cualquiera del cuerpo del texto con la paginación.
    _PRINTED_PAGE_PATTERNS = (
        re.compile(r'^p[aá]g(?:ina)?\.?\s*(\d{1,4})$', re.IGNORECASE),
        re.compile(r'^[-–—]?\s*(\d{1,4})\s*[-–—]?$'),
    )

    def _detect_printed_page_number(self, page_text: str) -> Optional[int]:
        """
        Busca el número de página impreso en la primera/última línea de una
        página. Devuelve None si no encuentra nada que calce con confianza
        (nunca "inventa" un número) — en ese caso, quien use esto debe caer
        de vuelta al índice físico de la página como aproximación.
        """
        lines = [l.strip() for l in page_text.split('\n') if l.strip()]
        if not lines:
            return None
        candidates = [lines[0], lines[-1]] if len(lines) > 1 else [lines[0]]
        for line in candidates:
            for pattern in self._PRINTED_PAGE_PATTERNS:
                match = pattern.match(line)
                if match:
                    number = int(match.group(1))
                    if 0 < number < 5000:  # sanity: descarta años, códigos, etc.
                        return number
        return None

    _PAGE_MARKER_PREFIX = '\x00P'
    _PAGE_MARKER_SUFFIX = '\x00'

    def _join_pages_with_markers(self, physical_pages: List[int], page_texts: List[str]) -> str:
        """
        Une el texto de cada página con '\\n\\n', igual que antes, pero
        antepone a cada una un marcador invisible (byte de control + número
        de página física, ej. "\\x00P23\\x00") que nadie ve en pantalla. Sirve
        para que el fragmentado por tokens (_split_into_chunks, en
        views_document_processor.py) pueda reconstruir de qué página física
        salió cada fragmento realmente enviado a la IA, y así citar la
        fuente exacta de cada pregunta generada en vez de "todo el capítulo".
        El marcador se descarta ahí antes de mandar el texto al modelo.
        """
        return '\n\n'.join(
            f'{self._PAGE_MARKER_PREFIX}{page}{self._PAGE_MARKER_SUFFIX}{text}'
            for page, text in zip(physical_pages, page_texts)
        )

    # ========================================================================
    # DOCX PROCESSING
    # ========================================================================
    
    def process_docx(self, file_path: str) -> Dict:
        """
        Procesa un archivo DOCX y extrae estructura por estilos (Heading 1, 2, etc.).
        
        Returns:
            Diccionario similar a process_pdf con capítulos organizados
        """
        result = {
            'metadata': {},
            'chapters': [],
            'full_text': '',
            'stats': {}
        }
        
        doc = Document(file_path)
        
        # Extraer metadata
        core_props = doc.core_properties
        result['metadata'] = {
            'title': core_props.title or '',
            'author': core_props.author or '',
            'subject': core_props.subject or '',
            'total_paragraphs': len(doc.paragraphs)
        }
        
        # Organizar por headings
        current_chapter = None
        chapters_list = []
        
        for para in doc.paragraphs:
            # Detectar si es un título (Heading 1)
            if para.style and para.style.name and para.style.name.startswith('Heading 1'):
                # Guardar capítulo anterior si existe
                if current_chapter:
                    chapters_list.append(current_chapter)
                
                # Iniciar nuevo capítulo
                current_chapter = {
                    'title': para.text,
                    'content': '',
                    'tokens': 0
                }
            elif current_chapter:
                # Añadir contenido al capítulo actual
                current_chapter['content'] += para.text + '\n\n'
        
        # Guardar último capítulo
        if current_chapter:
            chapters_list.append(current_chapter)

        # Descartar secciones sin contenido (dos Heading 1 seguidos, o el
        # último heading del documento sin texto debajo) — no tiene sentido
        # ofrecerlas como capítulo seleccionable para generar preguntas.
        chapters_list = [ch for ch in chapters_list if ch['content'].strip()]

        # Si no hay capítulos, considerar todo el documento
        if not chapters_list:
            full_text = '\n\n'.join([p.text for p in doc.paragraphs])
            chapters_list.append({
                'title': 'Documento completo',
                'content': full_text,
                'tokens': self.count_tokens(full_text)
            })
        else:
            # Calcular tokens
            for chapter in chapters_list:
                chapter['tokens'] = self.count_tokens(chapter['content'])
        
        result['chapters'] = chapters_list
        result['full_text'] = '\n\n'.join([ch['content'] for ch in chapters_list])
        result['stats'] = {
            'total_chapters': len(chapters_list),
            'total_tokens': sum(ch['tokens'] for ch in chapters_list)
        }
        
        return result
    
    # ========================================================================
    # PPTX PROCESSING
    # ========================================================================
    
    def process_pptx(self, file_path: str) -> Dict:
        """
        Procesa un archivo PowerPoint y extrae texto slide por slide.
        
        Returns:
            Diccionario con estructura por slides
        """
        result = {
            'metadata': {},
            'slides': [],
            'full_text': '',
            'stats': {}
        }
        
        prs = Presentation(file_path)
        
        result['metadata'] = {
            'title': prs.core_properties.title or '',
            'author': prs.core_properties.author or '',
            'total_slides': len(prs.slides)
        }
        
        slides_list = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_title = ''
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Primera shape con texto suele ser el título
                    if not slide_title and shape.text.strip():
                        slide_title = shape.text.strip()
                    else:
                        slide_text.append(shape.text)
            
            content = '\n'.join(slide_text)
            
            slides_list.append({
                'slide_number': i,
                'title': slide_title or f'Slide {i}',
                'content': content,
                'tokens': self.count_tokens(content)
            })
        
        result['slides'] = slides_list
        result['full_text'] = '\n\n'.join([s['content'] for s in slides_list])
        result['stats'] = {
            'total_slides': len(slides_list),
            'total_tokens': sum(s['tokens'] for s in slides_list)
        }

        # 'chapters' es el formato que consumen el generador de preguntas y el
        # dashboard (uno por slide); sin esta clave, cualquier PPTX queda con
        # 0 capítulos aunque la extracción haya funcionado.
        result['chapters'] = [
            {
                'title': s['title'],
                'content': s['content'] or s['title'],
                'tokens': s['tokens'] or self.count_tokens(s['title']),
                'pages': [s['slide_number']],
            }
            for s in slides_list if (s['content'] or s['title'])
        ]

        return result
    
    # ========================================================================
    # TOKEN COUNTING & OPTIMIZATION
    # ========================================================================
    
    def count_tokens(self, text: str) -> int:
        """
        Cuenta tokens exactos usando tiktoken (GPT-4 encoding).
        
        Args:
            text: Texto a contar
        
        Returns:
            Número de tokens
        """
        return len(self.encoding.encode(text))
    
    def split_by_token_limit(self, text: str, max_tokens: int = 4000) -> List[str]:
        """
        Divide texto en chunks que no excedan el límite de tokens.
        Útil para enviar a IA con ventana de contexto limitada.
        
        Args:
            text: Texto a dividir
            max_tokens: Límite de tokens por chunk
        
        Returns:
            Lista de chunks de texto
        """
        # Tokenizar
        tokens = self.encoding.encode(text)
        
        # Dividir en chunks
        chunks = []
        current_chunk = []
        
        for token in tokens:
            current_chunk.append(token)
            
            if len(current_chunk) >= max_tokens:
                # Decodificar chunk actual
                chunk_text = self.encoding.decode(current_chunk)
                chunks.append(chunk_text)
                current_chunk = []
        
        # Añadir último chunk si quedó algo
        if current_chunk:
            chunks.append(self.encoding.decode(current_chunk))
        
        return chunks
    
    def optimize_for_ai(self, text: str, 
                       remove_extra_whitespace: bool = True,
                       remove_urls: bool = False,
                       remove_emails: bool = False) -> str:
        """
        Optimiza texto para reducir tokens sin perder información relevante.
        
        Args:
            text: Texto a optimizar
            remove_extra_whitespace: Elimina espacios/saltos de línea extras
            remove_urls: Elimina URLs
            remove_emails: Elimina emails
        
        Returns:
            Texto optimizado
        """
        optimized = text
        
        if remove_extra_whitespace:
            # Eliminar múltiples espacios
            optimized = re.sub(r' +', ' ', optimized)
            # Eliminar múltiples saltos de línea (dejar máximo 2)
            optimized = re.sub(r'\n{3,}', '\n\n', optimized)
        
        if remove_urls:
            optimized = re.sub(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', '', optimized)
        
        if remove_emails:
            optimized = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '', optimized)
        
        return optimized.strip()
    
    # ========================================================================
    # EXPORT TO JSON
    # ========================================================================
    
    def export_to_json(self, result: Dict, output_file: str):
        """
        Exporta el resultado procesado a JSON.
        
        Args:
            result: Diccionario resultado de process_pdf/docx/pptx
            output_file: Ruta del archivo JSON de salida
        """
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
    
    def get_stats_summary(self, result: Dict) -> str:
        """
        Genera un resumen legible de las estadísticas del documento.
        
        Returns:
            String con resumen formateado
        """
        stats = result.get('stats', {})
        metadata = result.get('metadata', {})
        
        summary = []
        summary.append("=" * 60)
        summary.append("RESUMEN DEL DOCUMENTO")
        summary.append("=" * 60)
        
        if metadata.get('title'):
            summary.append(f"Título: {metadata['title']}")
        if metadata.get('author'):
            summary.append(f"Autor: {metadata['author']}")
        
        summary.append("")
        summary.append(f"Total de páginas/slides: {stats.get('total_pages', stats.get('total_slides', 'N/A'))}")
        summary.append(f"Total de capítulos/secciones: {stats.get('total_chapters', len(result.get('chapters', [])))}")
        summary.append(f"Total de tokens (GPT-4): {stats.get('total_tokens', 0):,}")
        
        if stats.get('removed_headers', 0) > 0:
            summary.append(f"Headers eliminados: {stats['removed_headers']}")
        if stats.get('removed_footers', 0) > 0:
            summary.append(f"Footers eliminados: {stats['removed_footers']}")
        
        summary.append("=" * 60)
        
        return '\n'.join(summary)


# ============================================================================
# FUNCIONES DE CONVENIENCIA
# ============================================================================

def process_document(file_path: str, **kwargs) -> Dict:
    """
    Función de conveniencia que detecta el tipo de archivo y procesa automáticamente.
    
    Args:
        file_path: Ruta al archivo (PDF, DOCX, PPTX)
        **kwargs: Argumentos adicionales pasados al procesador específico
    
    Returns:
        Diccionario con estructura procesada
    """
    processor = DocumentProcessor()
    
    if file_path.lower().endswith('.pdf'):
        return processor.process_pdf(file_path, **kwargs)
    elif file_path.lower().endswith('.docx'):
        return processor.process_docx(file_path, **kwargs)
    elif file_path.lower().endswith('.pptx'):
        return processor.process_pptx(file_path, **kwargs)
    else:
        raise ValueError(f"Formato no soportado: {file_path}")


def quick_token_count(file_path: str) -> int:
    """
    Cuenta rápidamente tokens de un documento sin procesamiento completo.
    
    Args:
        file_path: Ruta al archivo
    
    Returns:
        Número total de tokens
    """
    result = process_document(file_path)
    return result['stats'].get('total_tokens', 0)


# ============================================================================
# EJEMPLO DE USO
# ============================================================================

if __name__ == "__main__":
    # Ejemplo de uso básico
    processor = DocumentProcessor()
    
    # Procesar un PDF
    # result = processor.process_pdf("documento.pdf", remove_headers=True, remove_footers=True)
    
    # Ver resumen
    # print(processor.get_stats_summary(result))
    
    # Exportar a JSON
    # processor.export_to_json(result, "documento_procesado.json")
    
    # Contar tokens de un capítulo específico
    # capitulo1 = result['chapters'][0]
    # print(f"Capítulo: {capitulo1['title']}")
    # print(f"Tokens: {capitulo1['tokens']}")
    
    # Dividir por límite de tokens
    # chunks = processor.split_by_token_limit(capitulo1['content'], max_tokens=2000)
    # print(f"Dividido en {len(chunks)} chunks de máx 2000 tokens")
    
    print("Document Processor cargado correctamente.")
    print("Importa con: from document_processor import DocumentProcessor, process_document")
